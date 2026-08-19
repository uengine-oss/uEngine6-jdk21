# -*- coding: utf-8 -*-
"""
수협은행 BPM 도입 설명회 PPT 생성기
uEngine BPM  |  .sample/scenarios 기반
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ---------------------------------------------------------------- 브랜드 팔레트
# Sh수협은행 CI 계열(해양/블루) 를 참고한 팔레트
NAVY    = RGBColor(0x0A, 0x2A, 0x53)   # 딥 네이비 (본문 타이틀)
SH_BLUE = RGBColor(0x00, 0x60, 0xAE)   # 수협 블루 (주색)
SKY     = RGBColor(0x00, 0xA0, 0xD2)   # 아쿠아 (보조)
MINT    = RGBColor(0x00, 0xB3, 0xA4)   # 민트 (To-Be)
ORANGE  = RGBColor(0xF3, 0x7B, 0x21)   # 강조 (As-Is 문제)
RED     = RGBColor(0xD8, 0x3A, 0x3A)
GOLD    = RGBColor(0xE8, 0xB4, 0x2B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF4, 0xF8, 0xFC)   # 페이지 배경
CARD    = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xD8, 0xE4, 0xEF)
GRAY    = RGBColor(0x5A, 0x6B, 0x7D)
GRAY_L  = RGBColor(0x8B, 0x9A, 0xA8)
DARK    = RGBColor(0x1B, 0x2A, 0x3A)

FONT = "Pretendard"
FONT_FALLBACK = "Apple SD Gothic Neo"

SW, SH = 13.333, 7.5          # 16:9
M      = 0.72                 # 좌우 여백
CW     = SW - 2 * M           # 콘텐츠 폭

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_page = {"n": 0}

# ---------------------------------------------------------------- 기본 유틸
def _tf(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line=1.25, space_after=0, font=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tf

def textbox(slide, x, y, w, h, text, size=14, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.25, space_after=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _tf(tb, text, size, color, bold, align, anchor, line, space_after)
    return tb

def rect(slide, x, y, w, h, fill=CARD, outline=None, radius=None, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if outline is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = outline; shp.line.width = Pt(1.1)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp

def shape(slide, kind, x, y, w, h, fill=SH_BLUE, outline=None):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if outline is None: shp.line.fill.background()
    else: shp.line.color.rgb = outline; shp.line.width = Pt(1.1)
    shp.shadow.inherit = False
    return shp

def line_h(slide, x, y, w, color=LINE, weight=1.0):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln

def line_v(slide, x, y, h, color=LINE, weight=1.0):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln

# ---------------------------------------------------------------- 브랜드 마크
def brand_mark(slide, x, y, h=0.34, dark=False):
    """Sh 수협은행 워드마크 (설명회 자료용 텍스트 트리트먼트)"""
    box = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, h * 1.06, h,
                fill=(WHITE if dark else SH_BLUE))
    try: box.adjustments[0] = 0.24
    except Exception: pass
    _tf(box, "Sh", h * 46, (SH_BLUE if dark else WHITE), True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(slide, x + h * 1.24, y + h * 0.10, 1.4, h,
            "수협은행", size=h * 41, color=(WHITE if dark else NAVY), bold=True,
            anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- 슬라이드 프레임
def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    bgr = rect(s, 0, 0, SW, SH, fill=bg)
    bgr.shadow.inherit = False
    return s

def chrome(slide, kicker, title, sub=None, accent=SH_BLUE):
    """상단 타이틀 영역 + 하단 푸터. 콘텐츠 시작 y 를 돌려준다."""
    _page["n"] += 1
    # 상단 액센트 바
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, 0.075, fill=accent)
    if kicker:
        kw = 0.26 + sum((0.0155 if ord(c) > 0x2000 else 0.0078) * 11.0 for c in kicker)
        k = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, M, 0.34, kw, 0.28,
                  fill=RGBColor(0xE3, 0xEE, 0xF8))
        try: k.adjustments[0] = 0.5
        except Exception: pass
        _tf(k, kicker, 11, accent, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    ty = 0.68 if kicker else 0.46
    textbox(slide, M, ty, CW - 0.6, 0.58, title, size=26, color=NAVY, bold=True, line=1.1)
    y = ty + 0.64
    if sub:
        textbox(slide, M, y, CW - 0.6, 0.40, sub, size=13, color=GRAY, line=1.32)
        y += 0.30 + 0.17 * (sub.count("\n") + 1)
    line_h(slide, M, y - 0.08, CW, LINE, 1.0)
    footer(slide)
    return y + 0.12

def footer(slide):
    textbox(slide, M, SH - 0.46, 5.0, 0.26,
            "uEngine BPM  ·  수협은행 도입 설명회", size=9, color=GRAY_L)
    textbox(slide, SW - M - 1.2, SH - 0.46, 1.2, 0.26, f"{_page['n']:02d}",
            size=10, color=GRAY_L, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- 컴포넌트
def card(slide, x, y, w, h, title, body="", accent=SH_BLUE, num=None,
         title_size=14.5, body_size=11.5, fill=CARD, tcolor=None, bcolor=GRAY):
    box = rect(slide, x, y, w, h, fill=fill, outline=LINE, radius=0.055)
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.055, h, fill=accent)
    tx = x + 0.26
    cy = y + 0.20
    if num is not None:
        b = shape(slide, MSO_SHAPE.OVAL, tx, cy + 0.01, 0.30, 0.30, fill=accent)
        _tf(b, str(num), 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
        tx2 = tx + 0.40
    else:
        tx2 = tx
    textbox(slide, tx2, cy, w - (tx2 - x) - 0.24, 0.36, title,
            size=title_size, color=(tcolor or NAVY), bold=True, line=1.2)
    if body:
        nlines = title.count("\n") + 1
        textbox(slide, tx, cy + 0.34 * nlines + 0.12, w - 0.50, h - 0.62,
                body, size=body_size, color=bcolor, line=1.42)
    return box

def stat_tile(slide, x, y, w, h, value, label, accent=SH_BLUE, unit=""):
    rect(slide, x, y, w, h, fill=CARD, outline=LINE, radius=0.08)
    textbox(slide, x, y + h * 0.20, w, h * 0.42, value + unit,
            size=30, color=accent, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x, y + h * 0.66, w, h * 0.30, label,
            size=11, color=GRAY, align=PP_ALIGN.CENTER, line=1.25)

def chevron_flow(slide, x, y, w, h, items, fill=SH_BLUE, tsize=12, gap=0.06):
    n = len(items)
    cw = (w - gap * (n - 1)) / n
    for i, it in enumerate(items):
        cx = x + i * (cw + gap)
        sh_ = shape(slide, MSO_SHAPE.CHEVRON if i > 0 else MSO_SHAPE.PENTAGON,
                    cx, y, cw, h, fill=fill)
        _tf(sh_, it, tsize, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.2)

def arrow(slide, x, y, w, h, color=SH_BLUE, kind=MSO_SHAPE.RIGHT_ARROW):
    return shape(slide, kind, x, y, w, h, fill=color)

def badge(slide, x, y, text, color=SH_BLUE, size=10, h=0.26, pad=0.13):
    w = pad * 2 + sum((0.0145 if ord(c) > 0x2000 else 0.0072) * size for c in text)
    b = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=color)
    try: b.adjustments[0] = 0.5
    except Exception: pass
    _tf(b, text, size, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    return w

def bullets(slide, x, y, w, items, size=12.5, color=GRAY, gap=0.34,
            dot=SH_BLUE, bold_head=True):
    cy = y
    for it in items:
        if isinstance(it, tuple):
            head, tail = it
        else:
            head, tail = it, None
        d = shape(slide, MSO_SHAPE.OVAL, x, cy + 0.075, 0.075, 0.075, fill=dot)
        tb = slide.shapes.add_textbox(Inches(x + 0.20), Inches(cy - 0.03),
                                      Inches(w - 0.20), Inches(gap * 1.6))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.line_spacing = 1.35
        r = p.add_run(); r.text = head
        r.font.size = Pt(size); r.font.bold = bold_head
        r.font.color.rgb = (NAVY if bold_head and tail else color); r.font.name = FONT
        if tail:
            r2 = p.add_run(); r2.text = "  " + tail
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = FONT
        cy += gap
    return cy

def table(slide, x, y, w, rows, col_w, head_fill=NAVY, size=11, row_h=0.42,
          head_size=11, zebra=RGBColor(0xF2, 0xF7, 0xFB), align=None):
    """rows[0] 은 헤더. col_w 는 비율 리스트."""
    tot = sum(col_w)
    widths = [w * c / tot for c in col_w]
    cy = y
    for ri, row in enumerate(rows):
        cx = x
        rh = row_h if ri > 0 else row_h * 0.92
        if ri == 0:
            rect(slide, x, cy, w, rh, fill=head_fill, outline=None, radius=None)
        elif ri % 2 == 0:
            rect(slide, x, cy, w, rh, fill=zebra, outline=None)
        for ci, cell in enumerate(row):
            a = PP_ALIGN.LEFT if not align else align[ci]
            col = WHITE if ri == 0 else DARK
            tb = slide.shapes.add_textbox(Inches(cx + 0.14), Inches(cy),
                                          Inches(widths[ci] - 0.24), Inches(rh))
            _tf(tb, str(cell), head_size if ri == 0 else size, col,
                ri == 0, a, MSO_ANCHOR.MIDDLE, 1.18)
            cx += widths[ci]
        if ri > 0:
            line_h(slide, x, cy, w, LINE, 0.75)
        cy += rh
    line_h(slide, x, cy, w, LINE, 0.75)
    return cy

def section_divider(no, title, sub, items):
    s = new_slide(NAVY)
    _page["n"] += 1
    # 배경 장식
    c = shape(s, MSO_SHAPE.OVAL, SW - 3.6, -1.9, 5.4, 5.4, fill=RGBColor(0x10, 0x3A, 0x69))
    c2 = shape(s, MSO_SHAPE.OVAL, SW - 2.2, 3.4, 4.2, 4.2, fill=RGBColor(0x0D, 0x33, 0x5E))
    brand_mark(s, M, 0.5, 0.32, dark=True)
    textbox(s, M, 2.05, 1.6, 1.0, f"PART {no}", size=15, color=SKY, bold=True)
    textbox(s, M, 2.50, 8.6, 1.1, title, size=40, color=WHITE, bold=True, line=1.1)
    line_h(s, M, 3.62, 1.6, SKY, 3.0)
    textbox(s, M, 3.82, 8.0, 0.6, sub, size=14, color=RGBColor(0xB9, 0xD3, 0xE9), line=1.4)
    cy = 4.55
    for it in items:
        d = shape(s, MSO_SHAPE.OVAL, M + 0.02, cy + 0.085, 0.09, 0.09, fill=SKY)
        textbox(s, M + 0.28, cy - 0.02, 8.0, 0.34, it, size=13,
                color=RGBColor(0xD5, 0xE6, 0xF4))
        cy += 0.40
    textbox(s, SW - M - 1.2, SH - 0.46, 1.2, 0.26, f"{_page['n']:02d}",
            size=10, color=RGBColor(0x5C, 0x81, 0xA6), align=PP_ALIGN.RIGHT)
    return s

# ================================================================ 01 표지
s = new_slide(NAVY)
_page["n"] = 1
shape(s, MSO_SHAPE.OVAL, 8.6, -2.6, 7.2, 7.2, fill=RGBColor(0x10, 0x3B, 0x6B))
shape(s, MSO_SHAPE.OVAL, 10.4, 3.0, 5.4, 5.4, fill=RGBColor(0x0D, 0x33, 0x5D))
shape(s, MSO_SHAPE.OVAL, 9.9, 1.35, 2.9, 2.9, fill=SH_BLUE)
_tf(s.shapes[-1], "BPM", 30, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
brand_mark(s, M, 0.62, 0.36, dark=True)

textbox(s, M, 2.15, 8.4, 0.4, "업무 프로세스 관리(BPM) 도입 설명회", size=14, color=SKY, bold=True)
textbox(s, M, 2.66, 9.2, 1.9,
        "기록하는 시스템에서\n실행하는 시스템으로",
        size=44, color=WHITE, bold=True, line=1.18)
line_h(s, M, 4.72, 2.0, SKY, 3.5)
textbox(s, M, 4.98, 8.6, 1.0,
        "여신·수출환어음·신용카드·계좌개설·예금잔액통보 —\n수협 실제 업무 도면 7종으로 확인하는 BPM 적용 효과",
        size=14.5, color=RGBColor(0xC3, 0xDA, 0xEC), line=1.5)
line_h(s, M, 6.30, CW, RGBColor(0x1D, 0x4A, 0x7C), 1.0)
textbox(s, M, 6.48, 6.0, 0.4, "uEngine BPM  |  유엔진솔루션즈", size=12.5,
        color=RGBColor(0x9F, 0xC0, 0xDC))
textbox(s, SW - M - 4.0, 6.48, 4.0, 0.4, "설명회 자료 · 2026", size=12.5,
        color=RGBColor(0x9F, 0xC0, 0xDC), align=PP_ALIGN.RIGHT)

# ================================================================ 02 목차
s = new_slide()
y = chrome(s, "AGENDA", "오늘 말씀드릴 것",
           "BPM 을 처음 접하시는 분도 30분이면 «무엇이·왜·어떻게 달라지는가» 를 아실 수 있도록 구성했습니다.")
items = [
    ("1", "지금 우리 시스템의 한계", "계정계(SoR)만 있을 때 생기는 5가지 문제",
     "왜 «시스템은 다 있는데 일은 안 풀리는가»", SH_BLUE),
    ("2", "BPM 을 넣으면 무엇이 달라지나", "UI/UX · 준수율 · 변경속도 · 자동화 · 컴플라이언스",
     "찾아가는 업무에서 찾아오는 업무로", MINT),
    ("3", "수협 실제 업무로 보는 데모", "여신신규 외 7개 도면에서 도출한 8개 시연 시나리오",
     "우리 도면 그대로, 실행되는 프로세스로", SKY),
    ("4", "종합 · 기대효과 · 도입 로드맵", "무엇을 측정하고 어떻게 단계적으로 넓힐 것인가",
     "3단계 확산 계획과 성과 지표", NAVY),
]
cy = y + 0.16
for no, t, d, tag, col in items:
    rect(s, M, cy, CW, 1.10, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, M, cy, 0.06, 1.10, fill=col)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M + 0.34, cy + 0.27, 0.56, 0.56, fill=col)
    try: b.adjustments[0] = 0.25
    except Exception: pass
    _tf(b, no, 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, M + 1.10, cy + 0.22, 6.4, 0.36, t, size=17, color=NAVY, bold=True)
    textbox(s, M + 1.10, cy + 0.62, 7.4, 0.32, d, size=11.5, color=GRAY)
    textbox(s, SW - M - 4.5, cy + 0.40, 4.2, 0.36, tag, size=12,
            color=col, bold=True, align=PP_ALIGN.RIGHT)
    cy += 1.22

# ================================================================ PART 1
section_divider(1, "지금 시스템의 한계",
    "계정계·정보계는 «무엇이 일어났는가» 를 완벽히 기록합니다.\n그러나 «무엇을 해야 하는가» 는 어디에도 기록되어 있지 않습니다.",
    ["시스템 오브 레코드(SoR)가 하는 일과 하지 못하는 일",
     "사용자가 업무를 찾아다녀야 하는 구조",
     "문서와 사람 머릿속에만 있는 업무 프로세스",
     "정책이 바뀌면 개발이 따라가야 하는 구조",
     "수협 실제 도면 7종을 분석한 결과"])

# ================================================================ 03 SoR 정의
s = new_slide()
y = chrome(s, "PART 1", "시스템은 다 갖춰져 있는데, 왜 일은 안 풀릴까?",
           "은행 IT의 중심은 «기록의 시스템(System of Record)» 입니다. 기록은 완벽하지만, 기록만으로는 일이 흐르지 않습니다.")

rect(s, M, y + 0.10, CW * 0.485, 3.05, fill=CARD, outline=LINE, radius=0.05)
shape(s, MSO_SHAPE.RECTANGLE, M, y + 0.10, CW * 0.485, 0.06, fill=SH_BLUE)
textbox(s, M + 0.34, y + 0.34, 5.2, 0.34, "SoR — 기록의 시스템", size=17, color=NAVY, bold=True)
textbox(s, M + 0.34, y + 0.74, 5.2, 0.3, "계정계 · 여신원장 · 외환 · 카드 · 정보계",
        size=11.5, color=SH_BLUE, bold=True)
bullets(s, M + 0.34, y + 1.18, 5.1, [
    ("잘 하는 것", ""),
    ("· 계좌·거래·잔액을 정확히 기록", ""),
    ("· 단건 조회·수정·마감 처리", ""),
    ("· 정합성·감사증적 보존", ""),
], size=12, gap=0.30, dot=SH_BLUE, bold_head=False)
textbox(s, M + 0.34, y + 2.42, 5.1, 0.5,
        "그러나 «지금 이 건이 어느 단계에 있고,\n다음에 누가 무엇을 해야 하는가» 는 담기지 않습니다.",
        size=11.5, color=GRAY, line=1.35)

ax = M + CW * 0.515
rect(s, ax, y + 0.10, CW * 0.485, 3.05, fill=RGBColor(0xFF, 0xF6, 0xEE), outline=RGBColor(0xF6, 0xD9, 0xBC), radius=0.05)
shape(s, MSO_SHAPE.RECTANGLE, ax, y + 0.10, CW * 0.485, 0.06, fill=ORANGE)
textbox(s, ax + 0.34, y + 0.34, 5.2, 0.34, "빠져 있는 것 — 프로세스", size=17, color=NAVY, bold=True)
textbox(s, ax + 0.34, y + 0.74, 5.2, 0.3, "일의 순서 · 담당 · 기한 · 규칙 · 예외",
        size=11.5, color=ORANGE, bold=True)
bullets(s, ax + 0.34, y + 1.18, 5.1, [
    ("어디에 있나?", ""),
    ("· 업무 매뉴얼 PDF / 사내 규정집", ""),
    ("· 엑셀 체크리스트, 메신저 · 전화", ""),
    ("· 숙련 직원의 «머릿속»", ""),
], size=12, gap=0.30, dot=ORANGE, bold_head=False)
textbox(s, ax + 0.34, y + 2.42, 5.1, 0.5,
        "실행되지 않는 프로세스는 지켜지는지 알 수 없고,\n바꿔도 반영되었는지 확인할 수 없습니다.",
        size=11.5, color=GRAY, line=1.35)

qy = y + 3.42
rect(s, M, qy, CW, 0.78, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.40, qy + 0.16, CW - 0.8, 0.48,
        "“시스템은 무슨 일이 있었는지 알려주지만, 무슨 일을 해야 하는지는 알려주지 않는다.”",
        size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ 04 한계 5가지 개요
s = new_slide()
y = chrome(s, "PART 1", "SoR 만 있을 때 나타나는 5가지 한계",
           "수협 실제 업무 도면(여신신규 · 수출환어음 · 신용카드 · 계좌개설 · 예금잔액통보)에서 그대로 확인되는 문제들입니다.",
           accent=ORANGE)
data = [
    ("01", "찾아가는 업무", "사용자가 무엇을 할지\n직접 찾아 다녀야 한다",
     "여신 1건에 외부 사이트\n조회만 5회 이상"),
    ("02", "지켜지는지 모름", "프로세스가 문서에만 있어\n준수 여부를 측정할 수 없다",
     "게이트웨이에 조건식이 없고\nY/N 이름만 존재"),
    ("03", "느린 변경", "정책이 바뀌면 개발·배포가\n끝나야 현장에 반영된다",
     "우대금리 0.1% 조정에\n프로세스 재배포"),
    ("04", "자동화 공백", "틈새 업무는 여전히\n사람이 화면을 열어 처리",
     "시세·등기·진위확인이\n전부 수작업 UserTask"),
    ("05", "되돌릴 수 없음", "취소·정정 경로가 없어\n전화와 메모로 수습한다",
     "부결 시 이미 만든 보증서·\n담보등록 취소 경로 없음"),
]
cw = (CW - 0.24 * 4) / 5
for i, (no, t, d, ex) in enumerate(data):
    x = M + i * (cw + 0.24)
    rect(s, x, y + 0.14, cw, 3.95, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.14, cw, 0.075, fill=ORANGE)
    textbox(s, x + 0.22, y + 0.40, cw - 0.44, 0.4, no, size=26, color=RGBColor(0xF0, 0xC5, 0xA3), bold=True)
    textbox(s, x + 0.22, y + 0.92, cw - 0.44, 0.4, t, size=15.5, color=NAVY, bold=True)
    textbox(s, x + 0.22, y + 1.42, cw - 0.44, 1.0, d, size=11.5, color=GRAY, line=1.42)
    line_h(s, x + 0.22, y + 2.62, cw - 0.44, LINE, 1.0)
    textbox(s, x + 0.22, y + 2.78, cw - 0.44, 0.3, "우리 도면에서", size=9.5, color=GRAY_L, bold=True)
    textbox(s, x + 0.22, y + 3.04, cw - 0.44, 0.9, ex, size=10.5, color=ORANGE, line=1.42)

rect(s, M, y + 4.32, CW, 0.62, fill=RGBColor(0xE9, 0xF1, 0xF9), outline=None, radius=0.05)
textbox(s, M + 0.36, y + 4.44, CW - 0.7, 0.4,
        "다섯 가지 모두 «기록이 부족해서» 생긴 문제가 아닙니다. «실행할 프로세스가 시스템에 없어서» 생긴 문제입니다.",
        size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# ================================================================ 05 한계 ① 찾아가는 업무
s = new_slide()
y = chrome(s, "한계 ①", "사용자가 업무를 «찾아가야» 합니다",
           "계정계는 메뉴 구조입니다. 무엇을 해야 하는지 아는 사람만 메뉴를 열 수 있습니다.", accent=ORANGE)

textbox(s, M, y + 0.06, 5.0, 0.32, "AS-IS  ·  직원의 하루", size=12.5, color=ORANGE, bold=True)
rect(s, M, y + 0.44, CW * 0.485, 2.70, fill=RGBColor(0xFF, 0xF7, 0xF1), outline=RGBColor(0xF6, 0xDB, 0xC0), radius=0.05)
steps = [
    "출근한다. 오늘 내가 처리할 건이 몇 건인지 모른다.",
    "여신 메뉴를 연다 → 접수 대기 조회 → 없다.",
    "외환 메뉴를 연다 → 매입 대기 조회 → 3건 있다.",
    "메일함을 연다 → 지점장이 보완 요청한 건이 있다.",
    "엑셀 체크리스트를 연다 → 어제 놓친 건이 있다.",
    "결국 «내가 기억하는 만큼» 만 처리된다.",
]
cy = y + 0.66
for i, st in enumerate(steps):
    n = shape(s, MSO_SHAPE.OVAL, M + 0.28, cy + 0.015, 0.24, 0.24,
              fill=(RED if i == 5 else ORANGE))
    _tf(n, str(i + 1), 10, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, M + 0.64, cy, CW * 0.485 - 0.95, 0.3, st, size=11.5,
            color=(RED if i == 5 else DARK), bold=(i == 5))
    cy += 0.38

ax = M + CW * 0.515
textbox(s, ax, y + 0.06, 5.0, 0.32, "TO-BE  ·  BPM 워크리스트", size=12.5, color=MINT, bold=True)
rect(s, ax, y + 0.44, CW * 0.485, 2.70, fill=RGBColor(0xF0, 0xFB, 0xF9), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, ax + 0.28, y + 0.66, 5.0, 0.3, "내 할 일  (3)", size=13.5, color=NAVY, bold=True)
tasks = [
    ("여신신규 2026-0412 · 서류 검토", "오늘 마감", RED),
    ("수출환어음 EX-0088 · 매입 승인", "D-1", ORANGE),
    ("예금잔액통보 · 책임자 결재", "D-3", SH_BLUE),
]
cy = y + 1.00
for t, due, c in tasks:
    rect(s, ax + 0.28, cy, CW * 0.485 - 0.56, 0.46, fill=WHITE, outline=LINE, radius=0.08)
    shape(s, MSO_SHAPE.RECTANGLE, ax + 0.28, cy, 0.05, 0.46, fill=c)
    textbox(s, ax + 0.48, cy + 0.11, 3.9, 0.3, t, size=11, color=DARK)
    badge(s, ax + CW * 0.485 - 1.15, cy + 0.10, due, c, 9, 0.26)
    cy += 0.52
textbox(s, ax + 0.28, cy + 0.04, CW * 0.485 - 0.56, 0.56,
        "· 대기 중인 건은 자동으로 «배정»되어 목록에 올라옵니다.\n· 기한이 지나면 스스로 알리고, 상급자에게 상신됩니다.",
        size=10.5, color=GRAY, line=1.4)

by = y + 3.32
rect(s, M, by, CW, 0.64, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.40, by + 0.16, 2.2, 0.34, "UI / UX 의 전환", size=13, color=SKY, bold=True)
textbox(s, M + 2.70, by + 0.13, 9.8, 0.4,
        "메뉴를 «찾아가는» 화면  →  할 일이 «찾아오는» 화면 (Pull → Push)",
        size=16, color=WHITE, bold=True)
ry = by + 0.80
metrics = [
    ("업무 진입 경로", "메뉴 6~7단계 탐색", "워크리스트 1클릭"),
    ("누락 방지", "개인의 기억에 의존", "미처리 건 자동 잔류·알림"),
    ("교육 비용", "메뉴 위치를 외워야 함", "화면이 다음 할 일을 안내"),
]
mw = (CW - 0.24 * 2) / 3
for i, (k, a, b) in enumerate(metrics):
    x = M + i * (mw + 0.24)
    rect(s, x, ry, mw, 0.94, fill=CARD, outline=LINE, radius=0.06)
    textbox(s, x + 0.24, ry + 0.12, mw - 0.48, 0.3, k, size=11.5, color=NAVY, bold=True)
    textbox(s, x + 0.24, ry + 0.42, mw - 0.48, 0.3, "AS-IS   " + a, size=10.5, color=ORANGE)
    textbox(s, x + 0.24, ry + 0.66, mw - 0.48, 0.3, "TO-BE   " + b, size=10.5, color=MINT, bold=True)

# ================================================================ 06 한계 ② 준수율
s = new_slide()
y = chrome(s, "한계 ②", "프로세스가 «지켜졌는지» 알 수 없습니다",
           "그려 놓은 도면과 실제 실행 사이에 아무 연결이 없습니다. 도면은 그림이고, 실행은 사람의 판단입니다.", accent=ORANGE)

rect(s, M, y + 0.10, CW, 1.30, fill=CARD, outline=LINE, radius=0.05)
textbox(s, M + 0.32, y + 0.26, 4.0, 0.3, "수협 As-Is 도면에서 실제로 확인된 것", size=13, color=NAVY, bold=True)
ev = [("게이트웨이 20개", "조건식 0개 — 이름만 Y / N", RED),
      ("액티비티 약 120개", "입·출력 데이터 정의 0개", RED),
      ("프로세스 변수", "선언 0개 — 데이터가 흐르지 않음", RED),
      ("담당자 지정", "endpoint: null — 아무에게도 배정 안 됨", RED)]
ew = (CW - 0.64 - 0.20 * 3) / 4
for i, (k, v, c) in enumerate(ev):
    x = M + 0.32 + i * (ew + 0.20)
    textbox(s, x, y + 0.62, ew, 0.3, k, size=12.5, color=NAVY, bold=True)
    textbox(s, x, y + 0.92, ew, 0.4, v, size=10.5, color=c, line=1.3)

py = y + 1.58
textbox(s, M, py, 6.0, 0.32, "그 결과 — 준수율이 «측정 불가» 합니다", size=14, color=NAVY, bold=True)
qs = [
    ("이 건은 규정된 순서대로 처리됐나?", "확인 불가 — 순서를 강제하는 주체가 없음"),
    ("필수 서류 검토를 건너뛴 건은 몇 건인가?", "확인 불가 — 단계 이력이 남지 않음"),
    ("결재권한 기준을 넘겨 승인된 건이 있나?", "사후 감사로만 발견"),
    ("평균 처리시간이 가장 긴 단계는 어디인가?", "집계 불가 — 단계 개념 자체가 없음"),
]
cy = py + 0.44
for q, a in qs:
    rect(s, M, cy, CW * 0.52, 0.52, fill=RGBColor(0xF2, 0xF7, 0xFB), outline=None, radius=0.08)
    textbox(s, M + 0.24, cy + 0.13, CW * 0.52 - 0.4, 0.3, "Q.  " + q, size=12, color=NAVY, bold=True)
    arrow(s, M + CW * 0.53, cy + 0.15, 0.34, 0.22, RGBColor(0xC6, 0xD4, 0xE0))
    textbox(s, M + CW * 0.545 + 0.32, cy + 0.13, CW * 0.42, 0.3, a, size=12, color=RED)
    cy += 0.545

fy = cy + 0.14
rect(s, M, fy, CW, 0.76, fill=RGBColor(0xFF, 0xF3, 0xF3), outline=RGBColor(0xF3, 0xC9, 0xC9), radius=0.05)
textbox(s, M + 0.36, fy + 0.11, CW - 0.72, 0.56,
        "규제·감독 대응은 «했다» 가 아니라 «했음을 증명할 수 있다» 여야 합니다.\n실행되지 않는 도면은 증거가 되지 못합니다.",
        size=13.5, color=RED, bold=True, align=PP_ALIGN.CENTER, line=1.45)

# ================================================================ 07 한계 ③ 변경 속도
s = new_slide()
y = chrome(s, "한계 ③", "정책이 바뀌면, 개발이 끝나야 현장이 바뀝니다",
           "금융은 규칙이 자주 바뀌는 산업입니다. 그런데 그 규칙이 프로그램 코드 안에 들어 있습니다.", accent=ORANGE)

textbox(s, M, y + 0.06, 8.0, 0.32, "실제로 자주 바뀌는 규칙들 — 수협 여신신규 도면에서 추출", size=12.5, color=NAVY, bold=True)
rules = [
    ("DTI / DSR 산출", "정부 정책 변경 시 수시", RED),
    ("기준금리 · 가산금리 매트릭스", "월 / 수시", RED),
    ("조정 · 우대금리 적용", "수시", ORANGE),
    ("대출가능금액 (LTV)", "규제지역 지정 시", ORANGE),
    ("본부승인 대상 여부 (결재권한)", "반기", SH_BLUE),
    ("보증서 발급 가능 여부", "수시", ORANGE),
]
rw = (CW - 0.20 * 2) / 3
for i, (r, f, c) in enumerate(rules):
    x = M + (i % 3) * (rw + 0.20)
    yy = y + 0.44 + (i // 3) * 0.72
    rect(s, x, yy, rw, 0.62, fill=CARD, outline=LINE, radius=0.07)
    shape(s, MSO_SHAPE.RECTANGLE, x, yy, 0.05, 0.62, fill=c)
    textbox(s, x + 0.22, yy + 0.10, rw - 0.44, 0.3, r, size=12, color=NAVY, bold=True)
    textbox(s, x + 0.22, yy + 0.36, rw - 0.44, 0.26, "변경 빈도 · " + f, size=10, color=c)

cy = y + 2.06
textbox(s, M, cy, 6.0, 0.32, "AS-IS  ·  규칙 하나 바꾸는 데 걸리는 길", size=12.5, color=ORANGE, bold=True)
chevron_flow(s, M, cy + 0.36, CW, 0.62,
             ["정책 공고", "현업 요건 정리", "개발 요청", "코드 수정", "테스트", "배포 · 재기동", "현장 반영"],
             fill=RGBColor(0xB6, 0xC5, 0xD3), tsize=10.5)
textbox(s, M, cy + 1.06, CW, 0.3, "수 주 ~ 수 개월  ·  그 사이 현장은 «구 기준» 으로 처리합니다",
        size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)

cy2 = cy + 1.52
textbox(s, M, cy2, 6.0, 0.32, "TO-BE  ·  규칙을 프로세스 밖으로 꺼내면 (DMN)", size=12.5, color=MINT, bold=True)
chevron_flow(s, M, cy2 + 0.36, CW * 0.62, 0.62,
             ["정책 공고", "규칙표 수정", "저장", "다음 건부터 즉시 적용"], fill=MINT, tsize=12)
rect(s, M + CW * 0.645, cy2 + 0.36, CW * 0.355, 0.62, fill=NAVY, outline=None, radius=0.09)
textbox(s, M + CW * 0.645, cy2 + 0.50, CW * 0.355, 0.4,
        "프로세스 재배포 없음 · 무중단", size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ================================================================ 08 한계 ④ 자동화 공백
s = new_slide()
y = chrome(s, "한계 ④", "자동화되지 않은 «틈새»가 그대로 남아 있습니다",
           "큰 시스템은 자동화되어 있지만, 시스템과 시스템 사이의 일은 여전히 사람이 창을 열어 처리합니다.", accent=ORANGE)

textbox(s, M, y + 0.06, 9.0, 0.3,
        "수협 도면에서 «사람이 외부 화면을 여는» 것으로 그려진 업무", size=12.5, color=NAVY, bold=True)
rows = [
    ["원본 Task (도면 그대로)", "실제 대상 시스템", "현재", "BPM 적용 시"],
    ["KB / 국토교통부 시세조회", "KB부동산 · 실거래가", "사람이 웹 조회 후 수기 입력", "RPA Task — 봇이 조회 후 변수 자동 반영"],
    ["부동산 등기 결과 확인", "인터넷등기소", "사람이 발급 · 출력 · 스캔", "RPA Task — PDF 자동 첨부 + 항목 파싱"],
    ["공공마이데이터 정보 등록", "정부24", "사람이 조회 후 재입력", "RPA Task"],
    ["실명확인(신분증 진위확인)", "금융결제원", "사람이 화면 입력", "RPA Task / ServiceTask"],
    ["통보대상 계좌 조회 및 출력", "계정계 DB", "사람이 조회화면 실행", "SQL Task — 결과 N건이 프로세스 변수로"],
    ["예금잔액조회 기록부 보관", "계정계 DB", "사람이 등록화면 실행", "SQL Task (INSERT)"],
    ["수출대금 송금 전문 발송", "SWIFT 게이트웨이", "사람이 전문화면에서 전송", "Send Task — 전문 자동 송신"],
]
table(s, M, y + 0.42, CW, rows, [3.0, 2.2, 3.2, 4.0], head_fill=NAVY, size=11, row_h=0.44,
      align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])

by = y + 0.42 + 0.44 * 8 + 0.22
rect(s, M, by, CW * 0.485, 0.86, fill=RGBColor(0xFF, 0xF7, 0xF1), outline=RGBColor(0xF6, 0xDB, 0xC0), radius=0.05)
textbox(s, M + 0.28, by + 0.14, CW * 0.485 - 0.56, 0.6,
        "여신 1건 처리에 외부 사이트 조회만 5회 이상.\n조회한 값을 다시 손으로 옮겨 적으며 오타가 발생합니다.",
        size=11.5, color=DARK, line=1.4)
rect(s, M + CW * 0.515, by, CW * 0.485, 0.86, fill=RGBColor(0xF0, 0xFB, 0xF9), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, M + CW * 0.515 + 0.28, by + 0.14, CW * 0.485 - 0.56, 0.6,
        "BPM 은 봇을 «던져 놓고 잊는» 것이 아니라, 결과를 기다리고\n지연·실패를 SLA 로 관리하며 실패 시 사람으로 되돌립니다.",
        size=11.5, color=DARK, line=1.4)

# ================================================================ 09 한계 ⑤ 되돌림
s = new_slide()
y = chrome(s, "한계 ⑤", "이미 만들어 놓은 것을 «되돌릴» 방법이 없습니다",
           "여신은 앞 단계마다 계정계·보증기관·등기소에 실제 데이터를 남깁니다. 부결되면 그것을 전부 되돌려야 합니다.", accent=ORANGE)

textbox(s, M, y + 0.06, 9.0, 0.3, "여신신규 도면의 뒷부분 — 비가역적 등록의 연쇄", size=12.5, color=NAVY, bold=True)
seq = ["여신·담보\n예정등록", "보증서\n발급신청", "보증서 발급", "보증서\n담보등록", "여신 담보\n연결", "약정등록", "여신실행"]
cw = (CW - 0.14 * 6) / 7
for i, t in enumerate(seq):
    x = M + i * (cw + 0.14)
    rect(s, x, y + 0.44, cw, 0.76, fill=RGBColor(0xE8, 0xEF, 0xF6), outline=RGBColor(0xC9, 0xD9, 0xE7), radius=0.09)
    textbox(s, x, y + 0.56, cw, 0.6, t, size=10.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line=1.25)
    if i < 6:
        arrow(s, x + cw + 0.005, y + 0.72, 0.13, 0.18, RGBColor(0xB0, 0xC2, 0xD3))
textbox(s, M, y + 1.28, CW, 0.3,
        "각 단계가 외부 시스템에 실제 레코드를 남깁니다 — 계정계 / 서울보증 / 인터넷등기소",
        size=11, color=GRAY, align=PP_ALIGN.CENTER)

cy = y + 1.74
rect(s, M, cy, CW * 0.485, 2.30, fill=RGBColor(0xFF, 0xF6, 0xF6), outline=RGBColor(0xF3, 0xC9, 0xC9), radius=0.05)
badge(s, M + 0.28, cy + 0.20, "AS-IS", RED, 10)
textbox(s, M + 0.28, cy + 0.58, CW * 0.485 - 0.56, 0.3, "부결 통보하고 끝 — 뒤처리는 사람 몫", size=13.5, color=NAVY, bold=True)
bullets(s, M + 0.28, cy + 0.98, CW * 0.485 - 0.60, [
    "취소 경로가 도면에 아예 없다",
    "취소 순서(역순)가 강제되지 않는다",
    "어디까지 진행됐는지 사람이 기억해서 판단",
    "누락되면 유령 보증서 · 미해지 담보로 남는다",
], size=11.5, gap=0.32, dot=RED, bold_head=False)

ax = M + CW * 0.515
rect(s, ax, cy, CW * 0.485, 2.30, fill=RGBColor(0xF0, 0xFB, 0xF9), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
badge(s, ax + 0.28, cy + 0.20, "TO-BE · Compensation", MINT, 10)
textbox(s, ax + 0.28, cy + 0.58, CW * 0.485 - 0.56, 0.3, "완료된 것만, 역순으로, 엔진이 자동 취소", size=13.5, color=NAVY, bold=True)
bullets(s, ax + 0.28, cy + 0.98, CW * 0.485 - 0.60, [
    "«취소 화살표»를 한 개도 그리지 않는다",
    "2단계까지 갔으면 2개, 6단계면 6개만 취소",
    "담보연결 해제 → 담보등록 해제 → 보증서 반환 → 예정등록 취소",
    "되돌린 뒤 그 지점부터 «재개»까지 가능",
], size=11.5, gap=0.32, dot=MINT, bold_head=False)

fy = cy + 2.52
rect(s, M, fy, CW, 0.72, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.4, fy + 0.14, CW - 0.8, 0.44,
        "이 «자동 되돌림» 이 데모에서 가장 임팩트가 큰 장면입니다.  (시나리오 1 · 통합데모 12번)",
        size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ 10 As-Is 분석 데이터
s = new_slide()
y = chrome(s, "PART 1 · 진단", "수협 실제 도면 7종을 그대로 분석했습니다",
           "«BPM 이 필요하다» 는 주장이 아니라, 제출해 주신 도면에서 그대로 읽힌 숫자입니다.")

tiles = [("7", "분석한 업무 도면", SH_BLUE, "개"),
         ("약 120", "액티비티(Task)", SH_BLUE, ""),
         ("6", "사용된 BPMN 요소 종류", ORANGE, "종"),
         ("0", "프로세스 변수 선언", RED, "개"),
         ("0", "이벤트·자동화 Task 등", RED, "개")]
tw = (CW - 0.22 * 4) / 5
for i, (v, l, c, u) in enumerate(tiles):
    stat_tile(s, M + i * (tw + 0.22), y + 0.12, tw, 1.32, v, l, c, u)

cy = y + 1.66
textbox(s, M, cy, 6.0, 0.3, "도면에 한 번도 등장하지 않은 요소", size=12.5, color=NAVY, bold=True)
absent = ["Parallel / Inclusive / Event-based Gateway", "Boundary Event (Timer · Error · Message · Signal)",
          "Intermediate Catch / Throw Event", "SubProcess · Transaction · Event SubProcess",
          "Multi-Instance 마커", "Service / Send / Receive / Script / Business Rule Task",
          "Message Flow · Data Object · Data Store", "게이트웨이 조건식 (Y/N 이름만 존재)"]
aw = (CW - 0.18 * 3) / 4
for i, a in enumerate(absent):
    x = M + (i % 4) * (aw + 0.18)
    yy = cy + 0.38 + (i // 4) * 0.60
    rect(s, x, yy, aw, 0.52, fill=RGBColor(0xFA, 0xF0, 0xF0), outline=RGBColor(0xF0, 0xD5, 0xD5), radius=0.09)
    textbox(s, x + 0.14, yy + 0.09, aw - 0.28, 0.42, "✕  " + a, size=9.5, color=RED, line=1.25)

cy2 = cy + 1.68
textbox(s, M, cy2, 8.0, 0.3, "결론 — To-Be 로 가기 위한 4가지 정비 (전부 모델링 작업, 코드 개발 아님)",
        size=12.5, color=NAVY, bold=True)
fix = [("1", "프로세스 변수 선언", "여신번호 · 고객번호 · 금액 · 판정결과"),
       ("2", "게이트웨이 조건식 부여", "Y/N 이름 → 실제 조건식 + 기본 플로우"),
       ("3", "Task 타입 재지정", "UserTask → SQL · RPA · DMN · Send / Receive"),
       ("4", "계층 구조 도입", "보상 범위 · 반복 단위를 SubProcess 로")]
fw = (CW - 0.20 * 3) / 4
for i, (n, t, d) in enumerate(fix):
    x = M + i * (fw + 0.20)
    rect(s, x, cy2 + 0.34, fw, 1.06, fill=CARD, outline=LINE, radius=0.06)
    b = shape(s, MSO_SHAPE.OVAL, x + 0.22, cy2 + 0.48, 0.30, 0.30, fill=MINT)
    _tf(b, n, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, x + 0.62, cy2 + 0.48, fw - 0.84, 0.3, t, size=12.5, color=NAVY, bold=True)
    textbox(s, x + 0.22, cy2 + 0.88, fw - 0.44, 0.44, d, size=10.5, color=GRAY, line=1.3)

# ================================================================ PART 2
section_divider(2, "BPM 이 바꾸는 다섯 가지",
    "BPM 은 새 시스템을 하나 더 만드는 일이 아닙니다.\n이미 있는 시스템들 «위에» 업무의 흐름을 얹는 일입니다.",
    ["UI/UX — 찾아가는 업무에서 찾아오는 업무로",
     "준수율 — 그린 대로 실행되고, 그대로 측정된다",
     "변경 속도 — 그리면 즉시 조직 전체에 반영된다",
     "자동화 — RPA · SQL Task 로 틈새를 메운다",
     "컴플라이언스 — DMN 규칙 외부화와 설명 가능성"])

# ================================================================ 11 SoR + SoP
s = new_slide()
y = chrome(s, "PART 2", "기록의 시스템 위에 «실행의 시스템» 을 얹습니다",
           "계정계를 교체하지 않습니다. BPM 은 기존 시스템을 호출하고 조립하는 오케스트레이션 계층입니다.")

# 상단 : 사람
rect(s, M, y + 0.10, CW, 0.86, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.36, y + 0.24, 4.0, 0.3, "사용자 계층", size=11, color=SKY, bold=True)
uu = ["영업점 담당", "지점장", "본부 심사", "심사위원", "고객"]
uw = (CW - 5.2) / 5
for i, u in enumerate(uu):
    rect(s, M + 4.6 + i * (uw + 0.12), y + 0.24, uw, 0.46, fill=RGBColor(0x14, 0x44, 0x76), outline=None, radius=0.10)
    textbox(s, M + 4.6 + i * (uw + 0.12), y + 0.36, uw, 0.3, u, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# 중단 : BPM
my = y + 1.14
rect(s, M, my, CW, 1.72, fill=RGBColor(0xE7, 0xF4, 0xFB), outline=SKY, radius=0.05)
textbox(s, M + 0.36, my + 0.16, 6.0, 0.34, "uEngine BPM  —  실행의 시스템 (System of Process)",
        size=15, color=NAVY, bold=True)
mods = [("워크리스트 · 폼", "할 일이 찾아온다"), ("프로세스 엔진", "순서 · 담당 · 기한 강제"),
        ("DMN 룰 엔진", "규칙 외부화"), ("이벤트 · 타이머", "SLA · 예외 · 되돌림"),
        ("모니터링 · 통계", "준수율 · 병목 측정")]
mw = (CW - 0.72 - 0.16 * 4) / 5
for i, (t, d) in enumerate(mods):
    x = M + 0.36 + i * (mw + 0.16)
    rect(s, x, my + 0.62, mw, 0.92, fill=WHITE, outline=RGBColor(0xC5, 0xE0, 0xF0), radius=0.08)
    textbox(s, x + 0.14, my + 0.76, mw - 0.28, 0.3, t, size=11.5, color=SH_BLUE, bold=True,
            align=PP_ALIGN.CENTER)
    textbox(s, x + 0.14, my + 1.06, mw - 0.28, 0.4, d, size=10, color=GRAY,
            align=PP_ALIGN.CENTER, line=1.3)

# 하단 : SoR
by = my + 2.04
rect(s, M, by, CW, 1.30, fill=CARD, outline=LINE, radius=0.05)
textbox(s, M + 0.36, by + 0.16, 6.0, 0.3, "기록의 시스템 (System of Record) — 그대로 유지",
        size=13, color=NAVY, bold=True)
sors = ["계정계 · 여신원장", "외환 · SWIFT 게이트웨이", "카드 기간계", "정보계 · DW",
        "보증기관 · 등기소", "RPA 봇"]
sw_ = (CW - 0.72 - 0.14 * 5) / 6
for i, t in enumerate(sors):
    x = M + 0.36 + i * (sw_ + 0.14)
    rect(s, x, by + 0.58, sw_, 0.54, fill=RGBColor(0xEF, 0xF3, 0xF7), outline=LINE, radius=0.09)
    textbox(s, x + 0.08, by + 0.70, sw_ - 0.16, 0.32, t, size=10.5, color=DARK,
            align=PP_ALIGN.CENTER, line=1.25)

for xx in (M + 2.2, M + 6.4, M + 10.6):
    arrow(s, xx, my - 0.16, 0.20, 0.14, SKY, MSO_SHAPE.UP_ARROW)
    arrow(s, xx + 0.5, my + 1.76, 0.20, 0.24, SKY, MSO_SHAPE.DOWN_ARROW)

# ================================================================ 12 변화① UI/UX
s = new_slide()
y = chrome(s, "변화 ①", "UI / UX — 업무를 찾아가지 않고, 업무가 찾아옵니다", accent=MINT,
           sub="메뉴를 «아는 사람» 만 일할 수 있는 구조에서, 화면이 다음 할 일을 안내하는 구조로 바뀝니다.")

screens = [
    ("내 할 일 (Worklist)", "지금 내가 처리할 건만 기한순으로",
     ["역할·조직도 기반 자동 배정", "기한 임박 · 초과 색상 구분", "위임 · 대결 · 재배정"], SH_BLUE),
    ("업무 화면 (Form)", "단계마다 필요한 항목만 보이는 폼",
     ["이전 단계 데이터 자동 채움", "체크포인트 · 첨부 · 결재의견", "권한별 읽기/쓰기 제어"], SKY),
    ("진행 현황 (Tracking)", "이 건이 지금 어디에 있는지 한 화면",
     ["도면 위에 현재 위치 표시", "누가 · 언제 · 무엇을 했는지", "고객 문의에 즉시 답변"], MINT),
    ("관리자 모니터 (Monitor)", "지연·병목·부하를 실시간으로",
     ["단계별 평균 처리시간", "적체 건수 · 담당자 부하", "SLA 초과 자동 상신"], NAVY),
]
cwd = (CW - 0.24 * 3) / 4
for i, (t, d, li, c) in enumerate(screens):
    x = M + i * (cwd + 0.24)
    rect(s, x, y + 0.12, cwd, 3.30, fill=CARD, outline=LINE, radius=0.05)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.12, cwd, 0.075, fill=c)
    # 화면 목업
    rect(s, x + 0.24, y + 0.42, cwd - 0.48, 1.02, fill=RGBColor(0xF2, 0xF6, 0xFA), outline=LINE, radius=0.06)
    rect(s, x + 0.24, y + 0.42, cwd - 0.48, 0.22, fill=c, outline=None)
    for k in range(3):
        rect(s, x + 0.40, y + 0.76 + k * 0.20, cwd - 0.92, 0.11,
             fill=RGBColor(0xD5, 0xE1, 0xEC), outline=None)
    rect(s, x + cwd - 0.92, y + 0.76, 0.22, 0.11, fill=c, outline=None)
    textbox(s, x + 0.24, y + 1.58, cwd - 0.48, 0.32, t, size=13.5, color=NAVY, bold=True)
    textbox(s, x + 0.24, y + 1.92, cwd - 0.48, 0.42, d, size=10.5, color=c, line=1.35)
    cy = y + 2.38
    for it in li:
        d_ = shape(s, MSO_SHAPE.OVAL, x + 0.26, cy + 0.07, 0.07, 0.07, fill=c)
        textbox(s, x + 0.44, cy - 0.02, cwd - 0.66, 0.34, it, size=10.5, color=GRAY, line=1.3)
        cy += 0.31

by = y + 3.62
rect(s, M, by, CW, 1.02, fill=RGBColor(0xEF, 0xFA, 0xF8), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, M + 0.36, by + 0.14, 3.4, 0.3, "현업이 체감하는 변화", size=12, color=MINT, bold=True)
chg = ["메뉴 탐색 6~7단계 → 1클릭", "«내가 기억해야 할 일» → 시스템이 기억",
       "신입 교육 = 메뉴 암기 → 화면이 안내", "«이 건 어떻게 됐죠?» → 조회 한 번"]
cw2 = (CW - 4.0 - 0.16 * 3) / 4
for i, c_ in enumerate(chg):
    textbox(s, M + 3.9 + i * (cw2 + 0.16), by + 0.46, cw2, 0.44, "· " + c_, size=11,
            color=DARK, line=1.35)

# ================================================================ 13 변화② 준수율
s = new_slide()
y = chrome(s, "변화 ②", "프로세스 준수율 — «그린 대로» 실행됩니다", accent=MINT,
           sub="도면이 곧 실행 코드입니다. 사람이 순서를 건너뛸 수 없고, 건너뛴 흔적이 남습니다.")

mech = [("순서 강제", "선행 단계가 끝나야 다음 단계가 열린다", "시퀀스 플로우 · 토큰"),
        ("담당 강제", "권한 없는 사람의 워크리스트에 뜨지 않는다", "역할 · 조직도 매핑"),
        ("기한 강제", "기한이 지나면 알리고, 더 지나면 상신한다", "Timer Boundary · Escalation"),
        ("규칙 강제", "판정은 사람이 아니라 규칙표가 한다", "DMN Business Rule"),
        ("증적 자동", "누가·언제·무엇을·왜 했는지 전부 남는다", "인스턴스 이력 · 판정 사유")]
mw = (CW - 0.20 * 4) / 5
for i, (t, d, tech) in enumerate(mech):
    x = M + i * (mw + 0.20)
    rect(s, x, y + 0.12, mw, 1.62, fill=CARD, outline=LINE, radius=0.06)
    ic = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.22, y + 0.26, 0.36, 0.36, fill=MINT)
    try: ic.adjustments[0] = 0.25
    except Exception: pass
    _tf(ic, "✓", 14, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, x + 0.22, y + 0.70, mw - 0.44, 0.3, t, size=13, color=NAVY, bold=True)
    textbox(s, x + 0.22, y + 1.00, mw - 0.44, 0.5, d, size=10.5, color=GRAY, line=1.32)
    textbox(s, x + 0.22, y + 1.46, mw - 0.44, 0.26, tech, size=9.5, color=SH_BLUE)

cy = y + 1.82
textbox(s, M, cy, 8.0, 0.3, "그래서 이런 질문에 «숫자로» 답할 수 있게 됩니다", size=13, color=NAVY, bold=True)
rows = [
    ["관리 질문", "AS-IS", "BPM 도입 후"],
    ["규정 순서대로 처리된 비율은?", "측정 불가", "프로세스 준수율 (%) — 정의 대비 실행 경로 일치율"],
    ["단계별 평균 처리시간은?", "측정 불가", "액티비티별 소요시간 · 대기시간 자동 집계"],
    ["기한을 넘긴 건은 몇 건인가?", "사후 수기 집계", "SLA 초과 건 실시간 목록 + 자동 상신 이력"],
    ["부결 사유의 분포는?", "담당자 메모", "DMN 판정 사유(note) · 적중 규칙번호로 집계"],
    ["병목 구간은 어디인가?", "감으로 추정", "적체 건수 상위 액티비티 실시간 표시"],
]
table(s, M, cy + 0.32, CW, rows, [3.6, 2.6, 6.2], head_fill=NAVY, size=11, row_h=0.365)

fy = cy + 0.32 + 0.365 * 6 + 0.14
rect(s, M, fy, CW, 0.54, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.4, fy + 0.08, CW - 0.8, 0.42,
        "감독 대응 · 내부통제 관점 :  «했다» 를 넘어 «했음을 시스템이 증명한다» 로",
        size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ 14 변화③ 변경 속도
s = new_slide()
y = chrome(s, "변화 ③", "프로세스 변경 — 그리면, 조직 전체가 즉시 바뀝니다", accent=MINT,
           sub="변경의 성격에 따라 «누가 · 얼마나 빨리» 바꿀 수 있는지가 달라집니다. 대부분은 개발 없이 끝납니다.")

lv = [("규칙만 바뀔 때", "DSR 한도 40% → 30%\n우대금리 0.1% 조정", "룰(DMN) 표만 수정 → 저장",
       "현업 담당자", "즉시 · 무중단", MINT),
      ("흐름이 바뀔 때", "결재 단계 추가\n서류 검토 단계 분리", "모델러에서 도면 수정 → 신버전 배포",
       "프로세스 담당자", "당일", SKY),
      ("연계가 바뀔 때", "새 외부 시스템 연동\n전문 규격 변경", "Send / Receive Task 속성 수정",
       "IT 담당자", "수일", SH_BLUE),
      ("화면이 바뀔 때", "입력 항목 추가", "폼 편집 + 변수 추가",
       "프로세스 담당자", "당일", NAVY)]
lw = (CW - 0.22 * 3) / 4
for i, (t, ex, how, who, spd, c) in enumerate(lv):
    x = M + i * (lw + 0.22)
    rect(s, x, y + 0.12, lw, 2.66, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.12, lw, 0.075, fill=c)
    textbox(s, x + 0.24, y + 0.34, lw - 0.48, 0.3, t, size=14, color=NAVY, bold=True)
    rect(s, x + 0.24, y + 0.72, lw - 0.48, 0.62, fill=RGBColor(0xF3, 0xF7, 0xFB), outline=None, radius=0.08)
    textbox(s, x + 0.36, y + 0.80, lw - 0.72, 0.5, ex, size=10.5, color=GRAY, line=1.35)
    textbox(s, x + 0.24, y + 1.48, lw - 0.48, 0.44, how, size=11.5, color=DARK, bold=True, line=1.35)
    line_h(s, x + 0.24, y + 2.00, lw - 0.48, LINE, 1.0)
    textbox(s, x + 0.24, y + 2.10, lw - 0.48, 0.26, "주체   " + who, size=10.5, color=GRAY)
    textbox(s, x + 0.24, y + 2.38, lw - 0.48, 0.26, "반영   " + spd, size=11, color=c, bold=True)

cy = y + 3.02
textbox(s, M, cy, 9.0, 0.3, "«한 번 바꾸면 전 조직에 동시에 적용된다» 가 BPM 의 본질입니다",
        size=13.5, color=NAVY, bold=True)
comp = [("AS-IS", "지점마다 다른 엑셀 · 매뉴얼 버전. 공문을 돌려도 실제 적용은 제각각.", ORANGE),
        ("TO-BE", "새 버전을 배포하면 그 시점 이후 시작되는 모든 건에 동일 적용. 진행 중인 건은 기존 버전 유지.", MINT)]
cy2 = cy + 0.40
for t, d, c in comp:
    rect(s, M, cy2, CW, 0.62, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, M, cy2, 0.06, 0.62, fill=c)
    badge(s, M + 0.26, cy2 + 0.17, t, c, 10)
    textbox(s, M + 1.30, cy2 + 0.17, CW - 1.6, 0.32, d, size=12, color=DARK)
    cy2 += 0.72

# ================================================================ 15 변화④ 자동화
s = new_slide()
y = chrome(s, "변화 ④", "자동화 — 남은 틈새를 «프로세스가 관리하는 자동화» 로", accent=MINT,
           sub="시스템 개발 없이, 프로세스 도면 위에서 자동화 단계를 끼워 넣습니다.")

auto = [
    ("SQL Task", "DB 조회 · 적재를 사람 없이", SH_BLUE,
     ["SQL 을 몰라도 «컬럼 매핑»만으로 INSERT/SELECT 자동 생성",
      "조회 결과 N건이 그대로 프로세스 변수 · 다음 단계 반복 수가 됨",
      "프로세스 트랜잭션에 참여 — 실패 시 함께 롤백",
      "실행된 실제 SQL 이 인스턴스 이력에 남아 감사 추적 가능"],
     "예금잔액통보 대상 계좌 추출 · 여신 예정등록 · 기록부 적재"),
    ("RPA Task", "외부 사이트 조회를 봇에게", MINT,
     ["봇에 명령을 발행하고 프로세스는 대기 — 담당자 워크리스트는 비어 있음",
      "결과가 오면 시세 · 등기 항목이 변수에 자동 매핑",
      "10분 무응답이면 Timer 가 발화해 «수동 조회» 경로로 우회",
      "3회 실패 시 사람에게 폴백 — 봇 장애가 업무 중단이 되지 않음"],
     "KB/국토부 시세조회 · 인터넷등기소 · 공공마이데이터 · 진위확인"),
    ("Send / Receive", "전문·통보를 자동 송수신", SKY,
     ["SWIFT MT103 · MT999 전문 자동 발송, 접수번호를 변수에 저장",
      "입금 통보(MT910) 도착 시 상관키로 해당 건을 찾아 자동 재개",
      "이메일 통보는 트랜잭션 커밋 시 발송 — 롤백되면 메일도 안 나감",
      "프로세스 간 통신 — 여신신규 ↔ 여신심사 비동기 연결"],
     "송금전문 · 조회전문 · DM 발송/반송 · 고객 통보 · 심사 요청/결과"),
]
aw = (CW - 0.26 * 2) / 3
for i, (t, d, c, li, ex) in enumerate(auto):
    x = M + i * (aw + 0.26)
    rect(s, x, y + 0.12, aw, 4.10, fill=CARD, outline=LINE, radius=0.05)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.12, aw, 0.075, fill=c)
    textbox(s, x + 0.28, y + 0.36, aw - 0.56, 0.34, t, size=17, color=NAVY, bold=True)
    textbox(s, x + 0.28, y + 0.76, aw - 0.56, 0.3, d, size=11.5, color=c, bold=True)
    cy = y + 1.12
    for it in li:
        d_ = shape(s, MSO_SHAPE.OVAL, x + 0.30, cy + 0.075, 0.075, 0.075, fill=c)
        textbox(s, x + 0.50, cy - 0.03, aw - 0.80, 0.58, it, size=10.5, color=GRAY, line=1.36)
        cy += 0.56
    line_h(s, x + 0.28, y + 3.44, aw - 0.56, LINE, 1.0)
    textbox(s, x + 0.28, y + 3.54, aw - 0.56, 0.52, "수협 적용 지점\n" + ex, size=10, color=SH_BLUE, line=1.4)

by = y + 4.32
rect(s, M, by, CW, 0.62, fill=RGBColor(0xEF, 0xFA, 0xF8), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, M + 0.4, by + 0.09, CW - 0.8, 0.46,
        "핵심은 «자동화를 늘리는 것» 이 아니라 «자동화 실패를 프로세스가 책임지는 것» 입니다.  지연·오류 시 사람으로 돌아오는 길이 항상 도면에 있습니다.",
        size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ 16 변화⑤ 컴플라이언스
s = new_slide()
y = chrome(s, "변화 ⑤", "컴플라이언스 — 규칙을 «표» 로 꺼내고, 판단 근거를 남깁니다", accent=MINT,
           sub="DMN(Decision Model and Notation) 은 업무 규칙을 표로 표현하는 국제 표준입니다. 현업이 읽고 직접 고칠 수 있습니다.")

# 왼쪽 : DMN 표 예시
rect(s, M, y + 0.10, CW * 0.53, 3.24, fill=CARD, outline=LINE, radius=0.05)
textbox(s, M + 0.28, y + 0.28, 6.0, 0.3, "여신 시스템 심사 규칙표  ·  loanScreening", size=13, color=NAVY, bold=True)
rows = [["신용등급", "DSR", "LTV", "판정", "사유"],
        ["≤ 3", "≤ 40%", "≤ 70%", "APPROVE", "자동승인 기준 충족"],
        ["≤ 6", "≤ 40%", "—", "MANUAL", "본부 심사 대상"],
        ["—", "> 40%", "—", "REJECT", "DSR 한도 초과"],
        ["—", "—", "—", "REJECT", "기타 부적격"]]
table(s, M + 0.28, y + 0.62, CW * 0.53 - 0.56, rows, [1.1, 1.0, 1.0, 1.3, 2.4],
      head_fill=SH_BLUE, size=10.5, head_size=10.5, row_h=0.40)
textbox(s, M + 0.28, y + 2.72, CW * 0.53 - 0.56, 0.52,
        "위에서부터 처음 맞는 행이 적용됩니다(hit policy: FIRST).\n순서가 곧 우선순위이고, 마지막 행이 기본값입니다.",
        size=10.5, color=GRAY, line=1.4)

# 오른쪽 : 효과
ax = M + CW * 0.55
eff = [("규칙이 눈에 보인다", "규정집·코드·엑셀에 흩어진 기준이 한 표로 모입니다.", SH_BLUE),
       ("현업이 직접 바꾼다", "DSR 40 → 30 을 표에서 고치고 저장하면 다음 건부터 적용.", MINT),
       ("판단 근거가 남는다", "판정 결과와 «사유», 적중한 규칙 번호가 건별로 기록됩니다.", SKY),
       ("설명할 수 있다", "«왜 부결됐습니까» 에 근거를 제시할 수 있습니다 — 민원·감독 대응.", NAVY)]
cy = y + 0.10
for t, d, c in eff:
    rect(s, ax, cy, CW * 0.45, 0.74, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, ax, cy, 0.055, 0.74, fill=c)
    textbox(s, ax + 0.24, cy + 0.10, CW * 0.45 - 0.48, 0.3, t, size=13, color=NAVY, bold=True)
    textbox(s, ax + 0.24, cy + 0.40, CW * 0.45 - 0.48, 0.3, d, size=10.5, color=GRAY)
    cy += 0.84

fy = y + 3.50
rect(s, M, fy, CW, 1.04, fill=RGBColor(0xFF, 0xFB, 0xEE), outline=RGBColor(0xEE, 0xDA, 0xA8), radius=0.05)
textbox(s, M + 0.36, fy + 0.14, 3.0, 0.3, "금융권에서 특히 중요한 이유", size=12, color=GOLD, bold=True)
pts = ["규칙 변경 이력이 버전 관리된다", "누가 언제 규칙을 바꿨는지 남는다",
       "«설명 가능성» — 자동판정의 근거 제시", "규칙과 프로세스의 책임 분리"]
pw = (CW - 3.6 - 0.16 * 3) / 4
for i, p in enumerate(pts):
    textbox(s, M + 3.5 + i * (pw + 0.16), fy + 0.42, pw, 0.54, "· " + p, size=10.5, color=DARK, line=1.35)

# ================================================================ 17 Before / After 종합
s = new_slide()
y = chrome(s, "PART 2 · 요약", "한 장으로 보는 Before / After",
           "«시스템을 하나 더 만드는 것» 이 아니라 «일하는 방식» 이 바뀝니다.")
rows = [
    ["관점", "SoR 만 있을 때 (AS-IS)", "BPM 도입 후 (TO-BE)", "가능케 하는 기능"],
    ["업무 진입", "메뉴를 찾아가 조회한다", "할 일이 워크리스트로 찾아온다", "역할·조직도 기반 자동 배정"],
    ["업무 순서", "사람이 기억하고 판단한다", "엔진이 순서를 강제한다", "시퀀스 플로우 · 게이트웨이 조건식"],
    ["데이터", "단계 간 데이터가 끊긴다", "프로세스 변수로 끝까지 흐른다", "변수 선언 · 입출력 매핑"],
    ["기한 관리", "duration 은 장식, 초과해도 무반응", "기한 초과 시 알림 → 자동 상신", "Timer Boundary · Escalation"],
    ["대기 업무", "«확인» 업무가 워크리스트에 계속 쌓임", "엔진이 대기, 사람은 아무것도 안 한다", "Message Catch · Event-based Gateway"],
    ["병렬 업무", "위원 수만큼 Task 를 복사해야 한다", "인원 수만큼 자동 분기, 실행 중 증감 가능", "Multi-Instance (forEachRole)"],
    ["업무 규칙", "코드·엑셀·머릿속에 분산", "규칙표(DMN)로 외부화, 현업이 수정", "Business Rule Task"],
    ["규칙 변경", "개발 → 테스트 → 재배포 (수 주)", "표 수정 → 저장 → 즉시 적용 (무중단)", "룰 단독 배포"],
    ["외부 조회", "사람이 사이트 6~7개를 연다", "봇이 조회하고 결과가 변수에 채워진다", "RPA Task + 실패 폴백"],
    ["DB 작업", "사람이 조회·등록 화면을 실행", "SQL Task 가 트랜잭션 안에서 처리", "SQLTask · DatabaseMapping"],
    ["취소 · 정정", "되돌림 경로 없음, 전화와 메모로 수습", "완료된 것만 역순으로 자동 보상", "Compensation · backToHere"],
    ["측정 · 감사", "«했다» 는 주장만 남는다", "준수율·처리시간·판정사유가 데이터로", "인스턴스 이력 · 모니터링"],
]
table(s, M, y + 0.04, CW, rows, [1.5, 3.6, 3.9, 3.0], head_fill=NAVY, size=10.3,
      head_size=10.8, row_h=0.375)

# ================================================================ PART 3
section_divider(3, "수협 실제 업무로 보는 데모",
    "가상의 예제가 아닙니다. 제출해 주신 7개 업무 도면에서\n그대로 도출한 8개 시연 시나리오입니다.",
    ["여신신규(주택담보대출) · 여신심사",
     "수출환어음 매입 · 추심 결재",
     "기업신용카드 신규발급 · 일반계좌신규",
     "예금잔액 통보",
     "그리고 이 모두를 하나로 엮은 통합 데모"])

# ---------------------------------------------------------------- 시나리오 템플릿
def scenario_slide(no, feature, headline, sources, asis, steps, wow, checks, accent=SH_BLUE):
    s = new_slide()
    y = chrome(s, f"시나리오 {no}  ·  {feature}", headline, accent=accent)
    # 근거 도면 뱃지
    bx = M
    textbox(s, M, y + 0.00, 1.5, 0.28, "근거 도면", size=10, color=GRAY_L, bold=True)
    bx = M + 0.88
    for src in sources:
        bx += badge(s, bx, y - 0.02, src, RGBColor(0x6B, 0x82, 0x99), 9.5, 0.26) + 0.10

    top = y + 0.38
    CH = 3.28
    # 좌 : As-Is 한계
    lw = CW * 0.300
    rect(s, M, top, lw, CH, fill=RGBColor(0xFF, 0xF7, 0xF1), outline=RGBColor(0xF3, 0xD8, 0xBC), radius=0.05)
    badge(s, M + 0.24, top + 0.16, "AS-IS 의 한계", ORANGE, 10)
    cy = top + 0.58
    for a in asis:
        d_ = shape(s, MSO_SHAPE.OVAL, M + 0.26, cy + 0.07, 0.075, 0.075, fill=ORANGE)
        textbox(s, M + 0.46, cy - 0.04, lw - 0.72, 0.62, a, size=10.5, color=DARK, line=1.36)
        cy += 0.64

    # 우 : To-Be 단계
    rx = M + lw + 0.24
    rw = CW - lw - 0.24
    rect(s, rx, top, rw, CH, fill=CARD, outline=LINE, radius=0.05)
    badge(s, rx + 0.24, top + 0.16, "TO-BE 시연 흐름", accent, 10)
    n = len(steps)
    sw_ = (rw - 0.48 - 0.155 * (n - 1)) / n
    for i, (st, dt) in enumerate(steps):
        x = rx + 0.24 + i * (sw_ + 0.155)
        rect(s, x, top + 0.56, sw_, 1.92, fill=RGBColor(0xF4, 0xF8, 0xFC), outline=RGBColor(0xDF, 0xE9, 0xF2), radius=0.07)
        b = shape(s, MSO_SHAPE.OVAL, x + 0.12, top + 0.68, 0.25, 0.25, fill=accent)
        _tf(b, str(i + 1), 9.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
        textbox(s, x + 0.44, top + 0.67, sw_ - 0.56, 0.42, st, size=10.5, color=NAVY, bold=True, line=1.22)
        textbox(s, x + 0.13, top + 1.16, sw_ - 0.26, 1.24, dt, size=9.8, color=GRAY, line=1.36)
        if i < n - 1:
            arrow(s, x + sw_ + 0.005, top + 1.44, 0.145, 0.16, RGBColor(0xC0, 0xD2, 0xE2))
    # 시연 포인트
    rect(s, rx + 0.24, top + 2.58, rw - 0.48, 0.56, fill=NAVY, outline=None, radius=0.10)
    textbox(s, rx + 0.40, top + 2.63, rw - 0.80, 0.48, "★  " + wow, size=11, color=WHITE, bold=True, line=1.35)

    # 하단 : 확인 포인트
    by = top + CH + 0.14
    textbox(s, M, by, 5.0, 0.28, "데모에서 직접 확인하실 것", size=12, color=NAVY, bold=True)
    cwc = (CW - 0.20 * (len(checks) - 1)) / len(checks)
    for i, (k, v) in enumerate(checks):
        x = M + i * (cwc + 0.20)
        rect(s, x, by + 0.32, cwc, 0.92, fill=CARD, outline=LINE, radius=0.07)
        shape(s, MSO_SHAPE.RECTANGLE, x, by + 0.32, 0.05, 0.92, fill=MINT)
        textbox(s, x + 0.20, by + 0.42, cwc - 0.40, 0.28, k, size=10.8, color=NAVY, bold=True)
        textbox(s, x + 0.20, by + 0.70, cwc - 0.40, 0.46, v, size=9.8, color=GRAY, line=1.35)
    return s

# ================================================================ 18 데모 개요 맵
s = new_slide()
y = chrome(s, "PART 3", "데모 시나리오 맵 — 어느 도면에서 무엇을 보여 드리나",
           "7개 원본 도면 → 8개 시연 시나리오. 각 시나리오는 독립적으로도 시연 가능합니다.")

# 좌 : 원본 도면
textbox(s, M, y + 0.06, 3.0, 0.3, "수협 As-Is 도면", size=12, color=GRAY_L, bold=True)
docs = ["1. 여신신규_주택담보대출", "2. 수출환어음 매입", "3. 기업신용카드 신규발급",
        "4. 일반계좌신규", "5. 예금잔액 통보", "수출환어음 매입 및 추심결재", "여신심사"]
dw = CW * 0.235
for i, d in enumerate(docs):
    rect(s, M, y + 0.40 + i * 0.52, dw, 0.44, fill=CARD, outline=LINE, radius=0.08)
    textbox(s, M + 0.18, y + 0.50 + i * 0.52, dw - 0.36, 0.3, d, size=10.5, color=DARK)

# 우 : 시나리오
sx = M + dw + 1.20
textbox(s, sx, y + 0.06, 4.0, 0.3, "시연 시나리오 (uEngine 특징 기능)", size=12, color=SH_BLUE, bold=True)
scn = [("1", "보상 처리 Compensation", "부결·불일치 시 자동 되돌림", ORANGE),
       ("2", "멀티 인스턴스 Multi-Instance", "위원 수만큼 병렬 심의", SH_BLUE),
       ("3", "이벤트 처리 Event Handling", "입금 대기 · SLA 상신 · 예외 감시", SKY),
       ("4", "SQL Task", "조회 결과가 곧 프로세스 데이터", MINT),
       ("5", "RPA Task", "봇이 조회하고 프로세스가 기다린다", MINT),
       ("6", "메시지 송·수신", "SWIFT 전문 · 프로세스 간 통신", SKY),
       ("7", "DMN 비즈니스 룰", "규칙 무중단 변경", GOLD),
       ("8", "통합 데모", "여신신규 To-Be 한 흐름에 전부", NAVY)]
sw2 = CW - dw - 1.20
for i, (n, t, d, c) in enumerate(scn):
    yy = y + 0.40 + i * 0.455
    rect(s, sx, yy, sw2, 0.39, fill=CARD, outline=LINE, radius=0.09)
    shape(s, MSO_SHAPE.RECTANGLE, sx, yy, 0.05, 0.39, fill=c)
    b = shape(s, MSO_SHAPE.OVAL, sx + 0.16, yy + 0.065, 0.26, 0.26, fill=c)
    _tf(b, n, 10, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, sx + 0.54, yy + 0.075, 3.6, 0.28, t, size=11.5, color=NAVY, bold=True)
    textbox(s, sx + 4.20, yy + 0.075, sw2 - 4.4, 0.28, d, size=10.5, color=GRAY)

# 연결 화살표
arrow(s, M + dw + 0.34, y + 1.90, 0.52, 0.36, RGBColor(0xC5, 0xD7, 0xE6))
textbox(s, M + dw + 0.10, y + 2.34, 1.0, 0.3, "도출", size=10, color=GRAY_L, align=PP_ALIGN.CENTER)

# ================================================================ 19 시나리오 1 보상
scenario_slide(1, "보상 처리 (Compensation)",
    "부결되면, 이미 만들어 놓은 등록을 «자동으로» 되돌립니다",
    ["1. 여신신규_주택담보대출", "2. 수출환어음 매입"],
    ["취소 시 앞 단계로 돌아가는 경로가 도면에 아예 없다",
     "되돌림을 그리면 취소 대상 8단계 × 취소 Task = 도면이 폭발한다",
     "담보연결 해제 → 보증서 반환 → 예정등록 취소 의 «역순» 이 강제되지 않는다",
     "5단계까지만 진행된 건을 어디까지 취소해야 하는지 표현할 수 없다"],
    [("예정등록", "계정계에 여신·담보 예정 레코드 생성"),
     ("보증서 발급", "보증기관 연계 → 담보등록 → 여신 담보 연결"),
     ("본부 결재", "승인 여부 = N (부결)"),
     ("자동 보상", "완료된 등록만 역순으로 자동 취소"),
     ("부결 통보", "취소가 모두 끝난 뒤 고객 통보 후 종료")],
    "«취소 화살표» 를 한 개도 그리지 않았는데, 엔진이 «완료된 것만 · 역순으로» 되돌립니다",
    [("2단계까지 진행 후 부결", "보상 핸들러 2개만 실행"),
     ("6단계까지 진행 후 부결", "4개 실행 · 순서가 정확한 역순"),
     ("수출환어음 사후 하자", "«매입실행» 지점까지 되돌린 뒤 그 지점부터 재개"),
     ("모니터링 화면", "보상된 단계 상태가 Compensated 로 표시")],
    accent=ORANGE)

# ================================================================ 20 시나리오 2 멀티인스턴스
scenario_slide(2, "멀티 인스턴스 (Multiple Instances)",
    "여신심사위원회 — 위원 수만큼 «자동으로» 갈라집니다",
    ["여신심사", "5. 예금잔액 통보", "1. 여신신규_주택담보대출"],
    ["위원이 5명이면 Task 5개를 복사해 붙여야 한다",
     "위원 교체·증원 시 프로세스 정의를 다시 배포해야 한다",
     "«3번 위원이 반대» 같은 개별 결과가 데이터로 남지 않는다",
     "과반 찬성 · 전원 동의 판정이 사람 머릿속에만 있다"],
    [("위원 지정", "심사부 간사가 이번 안건의 위원 5명 지정"),
     ("자동 분기", "엔진이 심의 서브프로세스를 5개 생성"),
     ("병렬 심의", "각 위원 워크리스트에 본인 건 1개만 표시"),
     ("완료 대기", "전원 완료까지 상위 프로세스가 대기"),
     ("의결 집계", "찬반 집계 후 DMN 이 과반 여부 판정")],
    "정의는 하나인데 실행 시 인원 수만큼 갈라집니다 — 심의 도중 위원을 추가하면 인스턴스가 1개 더 생깁니다",
    [("위원 5명 지정", "서브 인스턴스 5개 생성 · 라벨에 위원 이름"),
     ("위원 3명으로 재실행", "정의 변경 없이 3개만 생성"),
     ("심의 중 위원 1명 추가", "기존 4개는 유지 + 1개 추가 생성"),
     ("4명만 처리한 상태", "상위 프로세스는 계속 대기")],
    accent=SH_BLUE)

# ================================================================ 21 시나리오 3 이벤트
scenario_slide(3, "이벤트 처리 (Event Handling)",
    "기다리는 일은 사람이 아니라 «엔진» 이 합니다",
    ["수출환어음 매입 및 추심결재", "여신심사", "2. 수출환어음 매입"],
    ["입금 대기를 «폴링 루프» 로 그려 담당자 워크리스트에 «확인» 업무가 계속 쌓인다",
     "모든 Task 의 duration 이 5로 박혀 있고 이를 소비하는 로직이 없다 — SLA 부재",
     "«언제든 도착하는» 하자·부도 통보를 표현할 수단이 없다",
     "고객 철회 · 시스템 장애 시 흐름을 끊는 경로가 없다"],
    [("전문 발송", "추심은행에 조회 전문 자동 송신"),
     ("3갈래 대기", "입금 통보 · 미입금 타임아웃 · 부도 시그널"),
     ("워크리스트 0건", "대기 중 사람이 할 일은 아무것도 없다"),
     ("이벤트 수신", "입금 통보 도착 → 금액·통화·입금일 자동 매핑"),
     ("예외 분기", "3일 미입금 → 독촉 전문 / 부도 → 채권 이관")],
    "기한이 지나면 스스로 독촉하고, 더 지나면 상급자에게 상신합니다 — 사람의 기억에 의존하지 않습니다",
    [("프로세스 시작 후 워크리스트", "비어 있음 — 사람 할 일 0건"),
     ("입금 통보 수신", "즉시 재개 + 전문 필드가 변수에 자동 매핑"),
     ("«필요서류 보완» 방치", "D+3 알림(Task 유지) → D+7 취소 + 지점장 상신"),
     ("연동 시스템 강제 오류", "프로세스가 죽지 않고 «수기 심사» 로 우회")],
    accent=SKY)

# ================================================================ 22 시나리오 4 SQL Task
scenario_slide(4, "SQL Task",
    "조회 결과가 곧 «프로세스 데이터» 가 됩니다",
    ["5. 예금잔액 통보", "1. 여신신규_주택담보대출"],
    ["조회·등록이 전부 UserTask — 사람이 화면을 열고 버튼을 누르는 것으로 모델링됨",
     "outputData 가 전부 비어 있어 조회 결과가 프로세스로 들어오지 않는다",
     "여러 등록 Task 가 하나의 트랜잭션인지 아닌지 알 수 없다",
     "어떤 SQL 이 언제 실행됐는지 인스턴스에 남지 않는다"],
    [("기준 입력", "본부 담당자가 기준일자·기준금액만 입력"),
     ("자동 조회", "SQL Task(SELECT)가 대상 계좌 N건 추출"),
     ("변수 적재", "계좌번호·예금주·잔액이 다중 변수로 채워짐"),
     ("건수 분기", "0건이면 즉시 종료 / N건이면 계좌별 반복"),
     ("이력 적재", "SQL Task(INSERT)가 기록부에 자동 적재")],
    "SQL 을 몰라도 «컬럼 매핑» 만으로 INSERT / SELECT 가 자동 생성됩니다 — 현업이 직접 DB 연동을 완성",
    [("기준금액 1,000만원 실행", "해당 계좌만 N건 조회되어 변수에 적재"),
     ("인스턴스 디버그 정보", "실제 실행된 SQL 전문이 그대로 표시"),
     ("컬럼 매핑만으로 등록", "생성된 SQL 확인 + 정상 INSERT"),
     ("INSERT 직후 강제 실패", "같은 트랜잭션이므로 함께 롤백")],
    accent=MINT)

# ================================================================ 23 시나리오 5 RPA
scenario_slide(5, "RPA Task",
    "봇이 조회하고, 프로세스가 «기다리고 책임집니다»",
    ["1. 여신신규_주택담보대출", "3. 기업신용카드 신규발급", "4. 일반계좌신규"],
    ["여신 1건당 외부 사이트 조회만 5회 이상 — 창을 6~7개 열어야 한다",
     "조회한 시세·등기·소유여부를 다시 손으로 입력하며 오타가 발생한다",
     "수 분~수 시간 걸리는 실제 대기가 duration:5 로만 표시된다",
     "사이트 점검 중이면 담당자가 «알아서 나중에» — 재시도 정책이 없다"],
    [("명령 발행", "엔진이 봇에 조회 명령을 보내고 대기 진입"),
     ("봇 실행", "KB부동산·국토부 실거래가를 순회하며 수집"),
     ("콜백 수신", "시세·출처·기준일이 변수에 자동 매핑"),
     ("지연 대응", "10분 무응답 시 Timer 발화 → 탁상감정 경로"),
     ("실패 폴백", "3회 실패 시 «수기 조회» UserTask 로 전환")],
    "«봇에게 던져두고 잊는 자동화» 가 아니라, 프로세스가 결과를 기다리고 지연·실패까지 관리합니다",
    [("RPA 단계 진입 시 워크리스트", "비어 있음 · 단계 상태는 Running"),
     ("봇 결과 발행", "프로세스 즉시 재개 + 시세 변수 채워짐"),
     ("무응답 시나리오", "Timer 발화 → 수동 조회 경로로 우회"),
     ("실패 콜백", "재시도 3회 후 사람에게 폴백")],
    accent=MINT)

# ================================================================ 24 시나리오 6 메시지
scenario_slide(6, "메시지 송·수신 (Send / Receive)",
    "전문과 통보를 «시스템이» 주고받습니다",
    ["2. 수출환어음 매입", "수출환어음 매입 및 추심결재", "5. 예금잔액 통보"],
    ["전문 발송·통보가 전부 UserTask — 사람이 전문화면에서 전송 버튼을 누른다",
     "«수신» 을 표현할 수단이 없어 폴링 루프로 그려져 있다",
     "여신신규 ↔ 여신심사 가 Call Activity 로만 연결되어 비동기 통보가 불가능하다",
     "도착한 전문이 어느 건인지 매칭할 «상관키» 개념이 도면에 없다"],
    [("전문 송신", "SWIFT MT103 송금 전문 자동 발송"),
     ("접수번호 저장", "응답의 전문접수번호가 변수에 저장"),
     ("수신 대기", "입금 통보(MT910)를 상관키로 대기"),
     ("자동 매칭", "도착 전문이 상관키로 해당 건을 찾아 재개"),
     ("프로세스 간 통신", "여신신규 ↔ 여신심사 비동기 연결")],
    "심사가 며칠 걸려도 여신신규는 다른 병렬 작업을 계속합니다 — 동기 호출(Call Activity)과의 결정적 차이",
    [("전문 발송", "게이트웨이에 도달 + 접수번호가 변수에 저장"),
     ("상관키를 틀리게 보냄", "어떤 건도 재개되지 않음 — 오매칭 방지 확인"),
     ("이메일 통보 후 강제 롤백", "메일이 발송되지 않음 (트랜잭션 큐)"),
     ("여신심사 완료 통보", "대기 중이던 여신신규가 자동 재개")],
    accent=SKY)

# ================================================================ 25 시나리오 7 DMN
scenario_slide(7, "DMN 비즈니스 룰",
    "DSR 한도 40% → 30%,  프로세스를 «건드리지 않고» 바꿉니다",
    ["1. 여신신규_주택담보대출", "3. 기업신용카드 신규발급", "여신심사"],
    ["본부승인대상여부 게이트웨이에 조건식이 없고 Y/N 플로우 «이름» 만 있다",
     "금리 매트릭스가 별도 엑셀·화면에 흩어져 있다",
     "우대금리 0.1% 조정에도 프로세스 버전을 올려 재배포해야 한다",
     "«왜 부결됐나» 에 답할 데이터가 남지 않는다"],
    [("입력 수집", "신용등급·소득·부채·LTV 가 변수에 채워짐"),
     ("규칙 호출", "DMN 결정표가 판정과 사유를 반환"),
     ("자동 분기", "APPROVE / MANUAL / REJECT 로 게이트웨이 분기"),
     ("규칙만 수정", "DSR 한도 표를 40 → 30 으로 고치고 저장"),
     ("즉시 적용", "재배포 없이 새 건부터 새 기준 적용")],
    "«규칙 변경 = 프로세스 재배포» 를 깨는 장면입니다 — 현업 담당자가 규칙표를 직접 고칩니다",
    [("신용 2등급 · DSR 35%", "APPROVE 로 자동 진행"),
     ("DSR 45%", "REJECT + 사유 «DSR 한도 초과» 기록"),
     ("룰 파일만 수정", "프로세스 재배포 없이 새 인스턴스부터 즉시 적용"),
     ("인스턴스 상세", "판정 입력값·결과·사유가 남아 부결 사유 설명 가능")],
    accent=GOLD)

# ================================================================ 26 통합 데모 흐름
TYPE_COLOR = {"USER": RGBColor(0x8B, 0x9A, 0xA8), "SQL": MINT, "RPA": SKY,
              "MSG": SH_BLUE, "DMN": GOLD, "CMP": ORANGE, "EVT": RGBColor(0x7A, 0x5A, 0xC6)}
TYPE_MARK = {"USER": "", "SQL": "■", "RPA": "▲", "MSG": "★", "DMN": "◇", "CMP": "◆", "EVT": "◎"}

s = new_slide()
y = chrome(s, "시나리오 8  ·  통합 데모", "여신신규(주택담보대출) To-Be — 7개 기능을 한 흐름에",
           "규칙이 심사하고(DMN) · 봇이 조회하고(RPA) · DB를 읽고 쓰고(SQL) · 위원회가 병렬 심의하고(MI) · 전문을 주고받으며(Message) · 기한을 넘기면 상신하고(Event) · 부결되면 되돌립니다(Compensation).")

# 범례
lg = [("SQL", "SQL Task"), ("RPA", "RPA Task"), ("MSG", "Message 송·수신"),
      ("DMN", "DMN 규칙"), ("CMP", "Compensation"), ("EVT", "Event / Timer")]
lx = M
for k, t in lg:
    c = TYPE_COLOR[k]
    sq = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, lx, y + 0.02, 0.24, 0.24, fill=c)
    try: sq.adjustments[0] = 0.3
    except Exception: pass
    _tf(sq, TYPE_MARK[k], 9, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, lx + 0.30, y + 0.03, 1.5, 0.24, t, size=10, color=GRAY)
    lx += 0.34 + 0.085 * len(t) + 0.45

def node(x, yy, w, h, label, kind="USER"):
    c = TYPE_COLOR[kind]
    fill = WHITE if kind == "USER" else RGBColor(
        min(255, c[0] // 1 + (255 - c[0]) * 88 // 100),
        min(255, c[1] + (255 - c[1]) * 88 // 100),
        min(255, c[2] + (255 - c[2]) * 88 // 100))
    rect(s, x, yy, w, h, fill=fill, outline=(LINE if kind == "USER" else c), radius=0.10)
    if kind != "USER":
        textbox(s, x + 0.08, yy + 0.04, 0.3, 0.2, TYPE_MARK[kind], size=8.5, color=c, bold=True)
    textbox(s, x + 0.06, yy + 0.16, w - 0.12, h - 0.20, label, size=9.8,
            color=(DARK if kind == "USER" else NAVY), bold=(kind != "USER"),
            align=PP_ALIGN.CENTER, line=1.22)

LX, LW = M, 1.05           # 좌측 구간 라벨
FX = M + LW + 0.16
FW = CW - LW - 0.16
RH, RG = 0.62, 0.20
ry = y + 0.44

def flow_row(label, nodes, ry, band=None, bandcolor=None):
    rect(s, LX, ry, LW, RH, fill=RGBColor(0xE9, 0xF0, 0xF7), outline=None, radius=0.12)
    textbox(s, LX + 0.06, ry + 0.14, LW - 0.12, RH - 0.16, label, size=10.5,
            color=NAVY, bold=True, align=PP_ALIGN.CENTER, line=1.2)
    n = len(nodes)
    gap = 0.26
    nw = (FW - gap * (n - 1)) / n
    for i, (lb, kd) in enumerate(nodes):
        x = FX + i * (nw + gap)
        node(x, ry, nw, RH, lb, kd)
        if i < n - 1:
            arrow(s, x + nw + 0.03, ry + RH / 2 - 0.07, 0.20, 0.14, RGBColor(0xBD, 0xCE, 0xDD))

flow_row("접수 · 조회", [("여신 신청", "USER"), ("필요서류 징구 · 스캔", "USER"),
                     ("고객 거래정보 조회", "SQL"), ("공공마이데이터 등록", "RPA"),
                     ("개인신용정보 동의", "USER")], ry)
ry += RH + RG
# 반복 표시
rect(s, FX - 0.06, ry - 0.05, FW + 0.12, RH + 0.10, fill=None,
     outline=RGBColor(0x9B, 0x7D, 0xD9), radius=0.06)
textbox(s, FX + FW - 3.2, ry - 0.26, 3.2, 0.24, "●  Multi-Instance · 담보물건 건수만큼 반복",
        size=9.5, color=RGBColor(0x7A, 0x5A, 0xC6), bold=True, align=PP_ALIGN.RIGHT)
flow_row("담보 평가", [("담보물건 등록", "USER"), ("KB/국토부 시세조회", "RPA"),
                    ("외부 감정평가 의뢰", "MSG"), ("감정 회신 수신", "MSG"),
                    ("담보가액 산출", "USER")], ry)
ry += RH + RG + 0.10
flow_row("심사 · 판정", [("대출가능금액 LTV", "DMN"), ("DTI / DSR 산출", "DMN"),
                     ("시스템 심사", "DMN"), ("여신심사 요청 · 결과", "MSG"),
                     ("금리 · 우대 · 결재권한", "DMN")], ry)
ry += RH + RG
flow_row("등록 · 실행", [("여신·담보 예정등록", "SQL"), ("보증서 발급", "MSG"),
                     ("보증서 담보등록", "SQL"), ("약정등록", "SQL"),
                     ("부동산 등기 확인", "RPA"), ("여신실행", "SQL")], ry)

# 보상 라인
ry += RH + 0.16
rect(s, LX, ry, CW, 0.60, fill=RGBColor(0xFF, 0xF5, 0xEC), outline=RGBColor(0xF3, 0xD5, 0xB6), radius=0.06)
textbox(s, LX + 0.18, ry + 0.16, 2.6, 0.3, "◆  부결 · 등기 불일치 시", size=10.5, color=ORANGE, bold=True)
cmp_steps = ["약정 취소", "여신 담보 연결 해제", "보증서 담보등록 해제", "보증서 반환 요청", "예정등록 취소"]
cx = LX + 2.90
cwn = (CW - 2.90 - 0.22 * 4) / 5
for i, t in enumerate(cmp_steps):
    x = cx + i * (cwn + 0.22)
    rect(s, x, ry + 0.10, cwn, 0.40, fill=WHITE, outline=ORANGE, radius=0.14)
    textbox(s, x + 0.04, ry + 0.19, cwn - 0.08, 0.26, t, size=9.5, color=ORANGE,
            bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + cwn + 0.02, ry + 0.23, 0.18, 0.14, ORANGE, MSO_SHAPE.LEFT_ARROW)

# 상시 감시
ry += 0.74
rect(s, LX, ry, CW, 0.50, fill=RGBColor(0xF5, 0xF2, 0xFC), outline=RGBColor(0xDA, 0xCF, 0xF0), radius=0.06)
textbox(s, LX + 0.18, ry + 0.12, 2.6, 0.3, "◎  상시 감시 (Event SubProcess)", size=10.5,
        color=RGBColor(0x7A, 0x5A, 0xC6), bold=True)
evs = ["고객 철회 → 전체 보상 후 중단", "서류보완 D+3 알림 · D+7 상신",
       "연동 오류 → 수기 심사 우회", "감정 회신 지연 → 재촉·교체"]
ex = LX + 2.90
ewn = (CW - 2.90 - 0.18 * 3) / 4
for i, t in enumerate(evs):
    textbox(s, ex + i * (ewn + 0.18), ry + 0.14, ewn, 0.3, "· " + t, size=9.8,
            color=RGBColor(0x5A, 0x45, 0x92))

# ================================================================ 27 시연 순서
s = new_slide()
y = chrome(s, "PART 3", "권장 시연 순서 (30분)",
           "시간이 부족하면 아래 «축소안» 순으로 줄이되, 12번 «자동 되돌림» 은 남기시길 권합니다.")
rows = [
    ["#", "시연 내용", "보여 드릴 화면", "소요"],
    ["1", "As-Is 도면 열기 — «전부 UserTask 와 분기뿐»", "모델러", "2분"],
    ["2", "To-Be 도면 열기 — 요소 타입 색상 대비", "모델러", "2분"],
    ["3", "정상 건 실행 : 신용 2등급 · DSR 35%", "인스턴스 모니터", "5분"],
    ["4", "└ SQL Task 결과가 프로세스 변수에 채워짐", "인스턴스 변수 + 실행된 SQL", "—"],
    ["5", "└ RPA 명령 발행 → 봇 콜백 → 시세 반영", "메시지 콘솔 + 인스턴스", "—"],
    ["6", "└ DMN 판정 결과와 «사유» 확인", "인스턴스 변수 (판정 사유)", "—"],
    ["7", "규칙 무중단 변경 : DSR 40 → 30 후 재실행", "룰 편집 화면 → 새 인스턴스", "4분"],
    ["8", "고액 건 실행 → 여신심사 프로세스 자동 기동", "프로세스 2개 동시 모니터", "5분"],
    ["9", "└ 위원 5명 지정 → 서브 인스턴스 5개 생성", "인스턴스 트리", "—"],
    ["10", "└ 심의 중 위원 1명 추가 → 인스턴스 1개 추가", "인스턴스 트리", "—"],
    ["11", "감정 회신 지연 → Timer 발화 → 재촉 전문", "이벤트 로그", "3분"],
    ["12", "부결 처리 → 보상 자동 실행  ★ 하이라이트", "Compensated 상태 + 취소 이력", "5분"],
    ["13", "등기 불일치 → 약정등록 지점으로 되돌린 후 재개", "인스턴스 모니터", "4분"],
]
table(s, M, y + 0.04, CW * 0.655, rows, [0.75, 5.0, 3.2, 0.9], head_fill=NAVY, size=10.5,
      head_size=10.5, row_h=0.345,
      align=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])

ax = M + CW * 0.675
textbox(s, ax, y + 0.04, 4.0, 0.3, "시간이 부족할 때 — 단계적 축소안", size=12.5, color=NAVY, bold=True)
red = [("최소 · 10분", "DMN + SQL Task + 게이트웨이 조건식", "선행 개발 없음", MINT),
       ("중간 · 20분", "위 + Multi-Instance + 메시지 송·수신 + Timer", "수신 엔드포인트 · 목서버 · 데모용 타이머", SKY),
       ("전체 · 30분", "위 + Compensation + RPA Task", "+ RPA 명령 발행부", SH_BLUE)]
cy = y + 0.44
for t, sc, pre, c in red:
    rect(s, ax, cy, CW * 0.325, 1.10, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, ax, cy, 0.055, 1.10, fill=c)
    textbox(s, ax + 0.24, cy + 0.12, CW * 0.325 - 0.48, 0.3, t, size=13, color=NAVY, bold=True)
    textbox(s, ax + 0.24, cy + 0.42, CW * 0.325 - 0.48, 0.42, sc, size=10.5, color=DARK, line=1.35)
    textbox(s, ax + 0.24, cy + 0.78, CW * 0.325 - 0.48, 0.3, "선행 : " + pre, size=10, color=GRAY)
    cy += 1.18

rect(s, ax, cy + 0.02, CW * 0.325, 0.98, fill=NAVY, outline=None, radius=0.06)
textbox(s, ax + 0.26, cy + 0.13, CW * 0.325 - 0.52, 0.76,
        "Compensation 은 선행 개발 없이 시연되고\n임팩트가 가장 큽니다 — 12번은 꼭 남기세요.",
        size=11.5, color=WHITE, bold=True, line=1.5)

# ================================================================ 28 데모 준비 현황
s = new_slide()
y = chrome(s, "PART 3", "데모 준비 현황 — 무엇이 되어 있고, 무엇을 준비해야 하나",
           "설명회 자리에서 «되는 것» 과 «준비 중인 것» 을 구분해 말씀드립니다.")
rows = [
    ["기능", "엔진 구현 상태", "데모 전 준비", "비고"],
    ["Compensation (보상)", "동작", "없음", "선행 개발 없이 즉시 시연 가능"],
    ["Multi-Instance (forEachRole)", "동작", "없음", "역할 다중 매핑만 설정"],
    ["Multi-Instance (forEachVariable)", "미구현 (스텁)", "선행 구현 또는 forEachRole 로 대체", "담보 건수 반복 시연 시에만 필요"],
    ["Event 처리 (Timer·Message·Signal·Error)", "동작", "데모용 타이머 단위 축소 버전", "72시간 → 분 단위"],
    ["SQL Task", "동작 (신규 개발 완료)", "데모 스키마 · 샘플 데이터", "DirectSQL / DatabaseMapping 모두"],
    ["RPA Task", "스텁 (명령 발행부 TODO)", "명령 발행 + 결과 수신 구현", "봇 오케스트레이터 REST 로 대체 가능"],
    ["Send / Receive Task", "동작 (POST·GET)", "전문 게이트웨이 목(mock) 서버", "PUT/PATCH 필요 시 확장"],
    ["DMN 비즈니스 룰", "동작 (Camunda DMN)", "룰 파일 작성", "JSON 결정표 · DMN XML 모두 지원"],
    ["외부 전문 수신 (상관키 → 재개)", "엔드포인트 구현 필요", "필수", "이벤트·메시지 시연의 전제"],
]
table(s, M, y + 0.04, CW, rows, [3.4, 2.8, 3.4, 3.4], head_fill=NAVY, size=10.5,
      head_size=11, row_h=0.40)

by = y + 0.04 + 0.40 * 10 + 0.20
rect(s, M, by, CW, 0.74, fill=RGBColor(0xEF, 0xFA, 0xF8), outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, M + 0.36, by + 0.10, CW - 0.72, 0.54,
        "«최소 10분 데모(DMN + SQL Task + 조건식)» 는 선행 개발 없이 오늘 상태로 시연 가능합니다.\n임팩트가 가장 큰 Compensation 역시 추가 개발 없이 시연됩니다.",
        size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line=1.5)

# ================================================================ PART 4
section_divider(4, "종합 · 기대효과 · 로드맵",
    "기능 이름이 아니라, 수협의 일이 어떻게 달라지는가로 정리합니다.",
    ["여섯 가지 효과 — 현업 · 관리 · IT 관점",
     "금융권에서 특히 의미가 큰 세 가지",
     "무엇을 측정해 성과를 확인할 것인가",
     "3단계 확산 로드맵",
     "왜 uEngine 인가"])

# ================================================================ 29 6가지 효과 종합
s = new_slide()
y = chrome(s, "PART 4 · 종합", "BPM 을 도입하면 결국 이 여섯 가지가 달라집니다",
           "앞서 보신 8개 시나리오를 «기능» 이 아니라 «효과» 로 다시 묶은 것입니다.")
eff = [
    ("01", "일이 사람을 찾아온다", ORANGE,
     "메뉴를 탐색하는 화면에서 할 일이 배정되는 화면으로. 누락은 시스템이 막고, 기한은 시스템이 챙깁니다.",
     "시나리오 3 · 워크리스트 0건"),
    ("02", "그린 대로 실행된다", SH_BLUE,
     "도면이 곧 실행입니다. 순서·담당·기한·규칙이 강제되고, 준수율을 숫자로 볼 수 있습니다.",
     "시나리오 2 · 3 · 준수율 측정"),
    ("03", "바꾸면 즉시 반영된다", MINT,
     "규칙은 표만 고치면 무중단 적용, 흐름은 도면을 고쳐 당일 배포. 전 지점에 동시에 적용됩니다.",
     "시나리오 7 · DSR 40 → 30"),
    ("04", "틈새가 자동화된다", SKY,
     "외부 사이트 조회는 봇이, DB 작업은 SQL Task 가. 실패하면 사람에게 되돌아오는 길이 도면에 있습니다.",
     "시나리오 4 · 5"),
    ("05", "규칙을 지키고 설명한다", GOLD,
     "판정 기준이 표로 드러나고, 판정 근거와 적중 규칙이 건별로 남습니다. 민원·감독 대응이 가능해집니다.",
     "시나리오 7 · 판정 사유 기록"),
    ("06", "되돌릴 수 있다", NAVY,
     "부결·오류 시 완료된 것만 역순으로 자동 취소. 유령 보증서·미해지 담보가 남지 않습니다.",
     "시나리오 1 · 자동 보상"),
]
cw3 = (CW - 0.24 * 2) / 3
for i, (n, t, c, d, tag) in enumerate(eff):
    x = M + (i % 3) * (cw3 + 0.24)
    yy = y + 0.08 + (i // 3) * 2.42
    rect(s, x, yy, cw3, 2.30, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, x, yy, cw3, 0.075, fill=c)
    textbox(s, x + 0.28, yy + 0.24, 1.0, 0.34, n, size=21, color=RGBColor(0xD5, 0xDF, 0xE8), bold=True)
    textbox(s, x + 0.28, yy + 0.68, cw3 - 0.56, 0.36, t, size=16, color=NAVY, bold=True)
    textbox(s, x + 0.28, yy + 1.08, cw3 - 0.56, 0.80, d, size=10.8, color=GRAY, line=1.42)
    line_h(s, x + 0.28, yy + 1.94, cw3 - 0.56, LINE, 1.0)
    textbox(s, x + 0.28, yy + 2.00, cw3 - 0.56, 0.26, tag, size=10, color=c, bold=True)

# ================================================================ 30 금융권 관점
s = new_slide()
y = chrome(s, "PART 4", "특히 «은행» 이기 때문에 의미가 큰 세 가지",
           "제조·유통과 달리, 금융은 규제·정합성·설명책임이 업무의 본질입니다.")
fin = [
    ("규제 대응", NAVY, "감독 대응이 «주장» 에서 «증거» 로",
     ["처리 순서·담당·시각이 건별로 자동 기록",
      "결재권한 기준을 시스템이 강제 — 초과 승인 원천 차단",
      "규칙 변경 이력이 버전으로 남아 «언제부터 이 기준» 인지 증명",
      "자동판정의 근거(사유·적중 규칙)를 제시 가능 — 설명 가능성"],
     "여신 결재권한 · DSR 한도 · 서류 충족 기준"),
    ("정합성 · 사고 예방", SH_BLUE, "«만들다 만 상태» 가 남지 않는다",
     ["부결 시 예정등록·보증서·담보연결을 역순 자동 취소",
      "SQL Task 가 프로세스 트랜잭션에 참여 — 실패 시 함께 롤백",
      "이메일·통보는 커밋 시점에 발송 — 롤백되면 나가지 않음",
      "중복 전문·중복 콜백에 대한 멱등 처리 지점이 명시됨"],
     "여신 등록 연쇄 · 수출환어음 사후 하자 · 대금 회수"),
    ("정책 대응 속도", MINT, "정부 정책이 바뀌면 «그 주에» 반영",
     ["DSR·LTV·규제지역 기준을 규칙표에서 즉시 수정",
      "금리·우대금리 매트릭스를 현업이 직접 관리",
      "프로세스 재배포·시스템 정지 없음",
      "진행 중인 건은 기존 기준 유지 — 소급 적용 사고 방지"],
     "DSR 한도 · LTV · 우대금리 · 수수료 감면"),
]
fw = (CW - 0.26 * 2) / 3
for i, (t, c, sub, li, ex) in enumerate(fin):
    x = M + i * (fw + 0.26)
    rect(s, x, y + 0.10, fw, 4.32, fill=CARD, outline=LINE, radius=0.05)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.10, fw, 0.075, fill=c)
    textbox(s, x + 0.28, y + 0.34, fw - 0.56, 0.36, t, size=18, color=NAVY, bold=True)
    textbox(s, x + 0.28, y + 0.78, fw - 0.56, 0.34, sub, size=12, color=c, bold=True, line=1.3)
    cy = y + 1.22
    for it in li:
        d_ = shape(s, MSO_SHAPE.OVAL, x + 0.30, cy + 0.075, 0.075, 0.075, fill=c)
        textbox(s, x + 0.50, cy - 0.03, fw - 0.80, 0.62, it, size=10.8, color=GRAY, line=1.38)
        cy += 0.60
    line_h(s, x + 0.28, y + 3.72, fw - 0.56, LINE, 1.0)
    textbox(s, x + 0.28, y + 3.81, fw - 0.56, 0.5, "수협 적용 지점\n" + ex, size=10, color=c, line=1.4)

rect(s, M, y + 4.52, CW, 0.56, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.4, y + 4.59, CW - 0.8, 0.42,
        "BPM 은 «편의 기능» 이 아니라 은행의 내부통제 인프라입니다.",
        size=14.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ================================================================ 31 성과 측정
s = new_slide()
y = chrome(s, "PART 4", "성과를 «어떻게 측정할 것인가»",
           "숫자를 먼저 약속드리지 않습니다. 지금은 «측정조차 안 되는» 지표들이 측정 가능해집니다 — 목표치는 진단 후 함께 정합니다.")

rows = [
    ["영역", "지표", "현재 측정", "BPM 도입 후 측정 방법"],
    ["속도", "건당 리드타임 (접수 → 실행)", "불가 — 기록 없음", "인스턴스 시작~종료 시각 자동 집계"],
    ["속도", "단계별 처리시간 · 대기시간", "불가", "액티비티별 소요/대기 시간 자동 집계"],
    ["품질", "프로세스 준수율", "불가", "정의 경로 대비 실행 경로 일치율"],
    ["품질", "기한(SLA) 초과 건수 · 비율", "수기 집계", "Timer 발화 이력 · 상신 이력"],
    ["품질", "재작업 · 되돌림 건수", "불가", "보상(Compensated) 실행 이력"],
    ["생산성", "1인당 동시 처리 건수", "불가", "워크리스트 배정·완료 건수"],
    ["생산성", "수작업 조회 · 재입력 횟수", "불가", "RPA / SQL Task 대체 단계 수"],
    ["통제", "결재권한 초과 승인 건수", "사후 감사", "DMN 판정 이력 + 권한 강제로 0 지향"],
    ["통제", "부결 사유 분포", "담당자 메모", "판정 사유(note) · 적중 규칙번호 집계"],
    ["대응", "정책 변경 반영 소요일", "수 주 ~ 수 개월", "규칙 수정 → 저장 시각으로 측정"],
]
table(s, M, y + 0.04, CW * 0.66, rows, [1.0, 3.6, 2.6, 4.0], head_fill=NAVY, size=10.3,
      head_size=10.5, row_h=0.375)

ax = M + CW * 0.685
textbox(s, ax, y + 0.04, 4.2, 0.3, "권장 접근", size=12.5, color=NAVY, bold=True)
step = [("1", "기준선(Baseline) 측정", "Phase 1 대상 업무를 BPM 에 올리고 3개월간 실측", SH_BLUE),
        ("2", "목표 합의", "실측값을 근거로 리드타임·준수율 목표를 함께 설정", SKY),
        ("3", "개선 반복", "병목 액티비티를 도면에서 바로 고쳐 재배포", MINT),
        ("4", "확산", "검증된 개선폭을 근거로 다음 업무군으로 확대", NAVY)]
cy = y + 0.44
for n, t, d, c in step:
    rect(s, ax, cy, CW * 0.315, 0.86, fill=CARD, outline=LINE, radius=0.06)
    b = shape(s, MSO_SHAPE.OVAL, ax + 0.22, cy + 0.28, 0.30, 0.30, fill=c)
    _tf(b, n, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, ax + 0.64, cy + 0.14, CW * 0.315 - 0.86, 0.3, t, size=12, color=NAVY, bold=True)
    textbox(s, ax + 0.64, cy + 0.44, CW * 0.315 - 0.86, 0.4, d, size=10, color=GRAY, line=1.35)
    cy += 0.96

rect(s, ax, cy + 0.06, CW * 0.315, 0.72, fill=RGBColor(0xFF, 0xFB, 0xEE),
     outline=RGBColor(0xEE, 0xDA, 0xA8), radius=0.06)
textbox(s, ax + 0.24, cy + 0.16, CW * 0.315 - 0.48, 0.52,
        "«측정할 수 없는 것은 개선할 수 없다»\nBPM 의 첫 번째 산출물은 데이터입니다.",
        size=11, color=NAVY, bold=True, line=1.45)

# ================================================================ 32 로드맵
s = new_slide()
y = chrome(s, "PART 4", "단계적 도입 로드맵",
           "전사를 한 번에 바꾸지 않습니다. 효과가 확실한 업무부터 올리고, 검증된 패턴을 확산합니다.")

ph = [("PHASE 1", "검증", "3~4개월", SH_BLUE,
       ["단순·정형 업무 1~2종을 먼저 (예: 일반계좌신규, 예금잔액 통보)",
        "워크리스트 · 폼 · 모니터링 정착",
        "SQL Task 로 계정계 조회·적재 연동",
        "DMN 으로 판정 규칙 1~2개 외부화"],
       "현업이 «찾아오는 업무» 를 체감하고 기준선 데이터를 확보"),
      ("PHASE 2", "확대", "6~9개월", SKY,
       ["복합 업무 확대 (수출환어음 매입 · 추심결재, 기업신용카드)",
        "메시지 송·수신으로 전문 연계 자동화",
        "Timer · Escalation 으로 SLA 관리 도입",
        "Multi-Instance 로 위원회·다건 처리 표준화"],
       "프로세스 준수율·리드타임이 지표로 관리되기 시작"),
      ("PHASE 3", "고도화", "9개월~", MINT,
       ["여신신규(주택담보대출) 등 최대 규모 업무 전환",
        "Compensation 으로 취소·정정 자동화",
        "RPA 연계로 외부 조회 전면 자동화",
        "프로세스 간 연결 · 전사 프로세스 지도 구축"],
       "규칙·흐름 변경을 현업이 주도하는 운영 체계로 전환")]
pw = (CW - 0.30 * 2) / 3
for i, (p, t, dur, c, li, out) in enumerate(ph):
    x = M + i * (pw + 0.30)
    rect(s, x, y + 0.34, pw, 3.96, fill=CARD, outline=LINE, radius=0.05)
    shape(s, MSO_SHAPE.RECTANGLE, x, y + 0.34, pw, 0.075, fill=c)
    bd = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.28, y + 0.10, 1.20, 0.44, fill=c)
    try: bd.adjustments[0] = 0.28
    except Exception: pass
    _tf(bd, p, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, x + 1.66, y + 0.16, pw - 1.9, 0.34, t, size=17, color=NAVY, bold=True)
    textbox(s, x + 0.28, y + 0.66, pw - 0.56, 0.3, "기간 · " + dur, size=11, color=c, bold=True)
    cy = y + 1.00
    for it in li:
        d_ = shape(s, MSO_SHAPE.OVAL, x + 0.30, cy + 0.075, 0.075, 0.075, fill=c)
        textbox(s, x + 0.50, cy - 0.03, pw - 0.80, 0.58, it, size=10.8, color=GRAY, line=1.38)
        cy += 0.56
    rect(s, x + 0.28, y + 3.30, pw - 0.56, 0.88, fill=RGBColor(0xF2, 0xF7, 0xFB), outline=None, radius=0.08)
    textbox(s, x + 0.42, y + 3.40, pw - 0.84, 0.70, "이 단계의 결과\n" + out, size=10.5,
            color=NAVY, line=1.4)
    if i < 2:
        arrow(s, x + pw + 0.05, y + 2.06, 0.20, 0.24, RGBColor(0xC5, 0xD7, 0xE6))

rect(s, M, y + 4.44, CW, 0.58, fill=RGBColor(0xEF, 0xFA, 0xF8),
     outline=RGBColor(0xBF, 0xE9, 0xE3), radius=0.05)
textbox(s, M + 0.4, y + 4.51, CW - 0.8, 0.42,
        "Phase 1 에서 만든 «변수 · 조건식 · 역할 · 룰» 은 Phase 2·3 에서 그대로 재사용됩니다 — 다시 만들지 않습니다.",
        size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# ================================================================ 33 왜 uEngine
s = new_slide()
y = chrome(s, "PART 4", "왜 uEngine BPM 인가",
           "표준을 지키면서, 실제 은행 업무에서 반드시 부딪히는 «예외» 를 다룰 수 있는가가 관건입니다.")
diff = [
    ("실행 중 인원 증감", "심의 도중 위원이 추가되면 서브 인스턴스가 1개 더 생깁니다. 대부분의 엔진은 분기 시점에 개수가 고정됩니다.",
     "refreshMultipleInstance()", SH_BLUE),
    ("지점 지정 되돌림", "«어디까지 되돌릴지» 를 지목하면 역순 보상 + 흐름 리셋 + 그 지점부터 재개까지 한 번에 처리합니다.",
     "backToHere() · compensateToThis()", ORANGE),
    ("SQL 없는 DB 연동", "현업이 컬럼 매핑만으로 SELECT/INSERT/UPDATE/DELETE 를 생성합니다. 폼 매핑과 동일한 화면을 씁니다.",
     "DatabaseMappingStrategy", MINT),
    ("규칙 단독 배포", "프로세스 정의를 건드리지 않고 규칙만 수정·배포합니다. JSON 결정표와 표준 DMN XML 을 모두 지원합니다.",
     "BusinessRuleTask + Camunda DMN", GOLD),
    ("트랜잭션 일관성", "SQL Task 가 프로세스 트랜잭션에 참여하고, 이메일은 커밋 시점에 발송됩니다. 롤백되면 아무것도 남지 않습니다.",
     "DataSourceConnectionFactory", SKY),
    ("표준 BPMN 2.0", "표준 도면을 그대로 읽고 실행합니다. 현업이 그린 As-Is 도면이 To-Be 의 출발점이 됩니다.",
     "BPMN 2.0 · DMN 1.3", NAVY),
]
dw = (CW - 0.24 * 2) / 3
for i, (t, d, tech, c) in enumerate(diff):
    x = M + (i % 3) * (dw + 0.24)
    yy = y + 0.12 + (i // 3) * 2.16
    rect(s, x, yy, dw, 1.92, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, x, yy, 0.055, 1.92, fill=c)
    textbox(s, x + 0.28, yy + 0.22, dw - 0.56, 0.34, t, size=15, color=NAVY, bold=True)
    textbox(s, x + 0.28, yy + 0.66, dw - 0.56, 0.86, d, size=11, color=GRAY, line=1.45)
    line_h(s, x + 0.28, yy + 1.54, dw - 0.56, LINE, 1.0)
    textbox(s, x + 0.28, yy + 1.62, dw - 0.56, 0.26, tech, size=9.5, color=c, bold=True)

rect(s, M, y + 4.44, CW, 0.60, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.4, y + 4.53, CW - 0.8, 0.42,
        "제출해 주신 7개 도면을 그대로 읽어 분석했습니다 — 새로 그리는 것이 아니라, 있는 도면을 «실행 가능하게» 만드는 일입니다.",
        size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ================================================================ 34 정리
s = new_slide()
y = chrome(s, "정리", "오늘 말씀드린 것을 한 문장으로",
           "BPM 은 새 시스템이 아니라, 이미 있는 시스템들을 «일의 순서» 로 꿰는 일입니다.")

rect(s, M, y + 0.14, CW, 1.24, fill=NAVY, outline=None, radius=0.05)
textbox(s, M + 0.5, y + 0.34, CW - 1.0, 0.84,
        "“계정계는 무슨 일이 있었는지 기록합니다.\nBPM 은 무슨 일을 해야 하는지 실행합니다.”",
        size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line=1.5)

cy = y + 1.62
sixp = [("찾아오는 업무", "메뉴 탐색 → 워크리스트", ORANGE),
        ("실행되는 도면", "그림 → 준수율 데이터", SH_BLUE),
        ("즉시 반영", "재배포 수 주 → 무중단 즉시", MINT),
        ("메워지는 틈새", "수기 조회 → RPA · SQL", SKY),
        ("설명되는 판단", "담당자 메모 → 판정 근거", GOLD),
        ("되돌릴 수 있음", "전화·메모 → 자동 역순 보상", NAVY)]
sw3 = (CW - 0.20 * 5) / 6
for i, (t, d, c) in enumerate(sixp):
    x = M + i * (sw3 + 0.20)
    rect(s, x, cy, sw3, 1.30, fill=CARD, outline=LINE, radius=0.06)
    shape(s, MSO_SHAPE.RECTANGLE, x, cy, sw3, 0.075, fill=c)
    textbox(s, x + 0.16, cy + 0.30, sw3 - 0.32, 0.34, t, size=13, color=NAVY, bold=True,
            align=PP_ALIGN.CENTER, line=1.2)
    textbox(s, x + 0.16, cy + 0.76, sw3 - 0.32, 0.44, d, size=10, color=GRAY,
            align=PP_ALIGN.CENTER, line=1.35)

ny = cy + 1.58
textbox(s, M, ny, 6.0, 0.3, "다음 단계 제안", size=13, color=NAVY, bold=True)
nx = [("1", "데모 일정 확정", "30분 통합 데모 · 축소안 중 선택"),
      ("2", "Phase 1 대상 업무 선정", "일반계좌신규 · 예금잔액 통보 등 정형 업무"),
      ("3", "기준선 측정 항목 합의", "리드타임 · 준수율 · SLA 초과 지표"),
      ("4", "PoC 범위·기간 확정", "3~4개월 · 산출물 정의")]
nw = (CW - 0.22 * 3) / 4
for i, (n, t, d) in enumerate(nx):
    x = M + i * (nw + 0.22)
    rect(s, x, ny + 0.36, nw, 1.10, fill=RGBColor(0xF2, 0xF7, 0xFB), outline=LINE, radius=0.06)
    b = shape(s, MSO_SHAPE.OVAL, x + 0.22, ny + 0.50, 0.30, 0.30, fill=SH_BLUE)
    _tf(b, n, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 1.0)
    textbox(s, x + 0.62, ny + 0.50, nw - 0.84, 0.3, t, size=12, color=NAVY, bold=True)
    textbox(s, x + 0.22, ny + 0.88, nw - 0.44, 0.46, d, size=10, color=GRAY, line=1.35)

# ================================================================ 35 Q&A
s = new_slide(NAVY)
_page["n"] += 1
shape(s, MSO_SHAPE.OVAL, 9.0, -2.4, 7.0, 7.0, fill=RGBColor(0x10, 0x3B, 0x6B))
shape(s, MSO_SHAPE.OVAL, 10.6, 3.2, 5.0, 5.0, fill=RGBColor(0x0D, 0x33, 0x5D))
brand_mark(s, M, 0.62, 0.36, dark=True)
textbox(s, M, 2.55, 8.0, 0.5, "감사합니다", size=44, color=WHITE, bold=True)
line_h(s, M, 3.55, 2.0, SKY, 3.5)
textbox(s, M, 3.82, 8.4, 1.2,
        "질문과 논의를 환영합니다.\n수협 업무 도면을 더 주시면 같은 방식으로 분석해 드리겠습니다.",
        size=15, color=RGBColor(0xC3, 0xDA, 0xEC), line=1.6)
rect(s, M, 5.30, 6.6, 1.02, fill=RGBColor(0x12, 0x3F, 0x70), outline=None, radius=0.06)
textbox(s, M + 0.36, 5.46, 6.0, 0.7,
        "uEngine BPM  ·  유엔진솔루션즈\nBPMN 2.0 / DMN 표준 기반 업무 프로세스 관리 플랫폼",
        size=12, color=RGBColor(0xA9, 0xC7, 0xE0), line=1.6)
textbox(s, SW - M - 1.2, SH - 0.46, 1.2, 0.26, f"{_page['n']:02d}",
        size=10, color=RGBColor(0x5C, 0x81, 0xA6), align=PP_ALIGN.RIGHT)

# ================================================================ 저장
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "수협은행_BPM_도입_설명회.pptx")
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
